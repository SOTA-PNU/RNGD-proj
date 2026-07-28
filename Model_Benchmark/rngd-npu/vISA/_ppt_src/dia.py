#!/usr/bin/env python3
"""손으로 그리는 도해 슬라이드들."""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from deck import (ACCENT, ACCENT_D, AMBER, AMBER_BG, BODY_BOT, BODY_W, CODE_BG, MR, E, vlen,
                  CODE_FG, GREEN, INK, LINE, ML, MUTED, RED, SOFT, SW, SH,
                  WHITE, arrow, hline, label, para, rect, style, tbox)

PALE = RGBColor(0xE3, 0xEE, 0xFB)
PALE2 = RGBColor(0xFD, 0xEB, 0xD3)
PALE3 = RGBColor(0xE2, 0xF4, 0xE8)
GREY = RGBColor(0xEC, 0xF0, 0xF4)


def _cap(s, x, y, w, text, size=11.5, color=MUTED, align=PP_ALIGN.LEFT, bold=False):
    """캡션. 상자 높이를 내용에 맞춰 잡는다(고정 높이면 아래 도형을 덮는다)."""
    per_line = max(4, int((w / 12700.0) / (size * 0.5)))   # 한 줄에 들어가는 반각 수
    lines = 0
    for seg in str(text).split("\n"):
        lines += max(1, -(-vlen(seg) // per_line))
    h = Pt(size * 1.28 * lines + 4)
    tf = tbox(s, x, y, w, h)
    para(tf, True, text, size, bold, color, align=align)


# ---------------------------------------------------------------- 1. 전체 지도
def roadmap(d):
    s = d._new()
    top = d._chrome(s, "이 발표의 전체 지도", "왼쪽에서 오른쪽으로 쌓아 올린다")
    cols = [
        ("1~2부\n기초", "행렬곱 · GEMM\neinsum · 텐서 축약\n메모리 벽 · 연산강도", PALE),
        ("3~4부\n왜 이렇게 만드나", "MAC · 시스톨릭\nTensor Core\n타일링", PALE2),
        ("5부\n하드웨어", "RNGD / TCP\nChip·Cluster·Slice·Lane\n8단계 파이프라인", PALE3),
        ("6~7부\nvISA 언어", "어느 계층인가\n매핑 m![]\n타일링을 적는 법", PALE),
        ("8~10부\n실제", "커널 해부\n스케줄 · 실측\n무엇이 되고 안 되나", GREY),
    ]
    n = len(cols)
    gap = Inches(0.16)
    w = (BODY_W - gap * (n - 1)) / n
    y = top + Inches(0.45)
    h = Inches(3.0)
    for i, (head, body, fill) in enumerate(cols):
        x = ML + i * (w + gap)
        box = rect(s, x, y, w, h, fill=fill, line=LINE)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.14)
        for j, ln in enumerate(head.split("\n")):
            para(tf, j == 0, ln, 15, True, ACCENT_D, align=PP_ALIGN.CENTER, space_after=1)
        for ln in body.split("\n"):
            para(tf, False, ln, 12, False, INK, align=PP_ALIGN.CENTER, space_before=3)
        if i < n - 1:
            arrow(s, x + w + Inches(0.01), y + h / 2, x + w + gap - Inches(0.01), y + h / 2)
    _cap(s, ML, y + h + Inches(0.28), BODY_W,
         "앞부분(1~4부)은 NPU가 아니어도 통하는 일반 지식이다. 5부부터 RNGD 고유 내용이 시작된다.",
         12.5, MUTED)
    d._callout(s, "지금 이해가 안 되는 부분이 있어도 좋다. 7부의 매핑 m![] 한 줄에 앞의 모든 개념이 모인다.")
    return s


# ---------------------------------------------------------------- 2. GEMM 타일링
def gemm_tiling(d):
    s = d._new()
    top = d._chrome(s, "타일링 그림으로 보기", "C 타일 하나 = A 행띠 × B 열띠")
    cell = Inches(0.34)
    N = 6
    gy = top + Inches(0.55)
    # A
    ax = ML + Inches(0.35)
    for i in range(N):
        for j in range(N):
            f = PALE if i == 2 else WHITE
            rect(s, ax + j * cell, gy + i * cell, cell, cell, fill=f, line=LINE, lw=0.75)
    _cap(s, ax, gy - Inches(0.32), cell * N, "A  (M×K)", 13, ACCENT_D, PP_ALIGN.CENTER, True)
    _cap(s, ax, gy + cell * N + Inches(0.08), cell * N, "행 타일 하나를 통째로 읽는다", 10.5, MUTED, PP_ALIGN.CENTER)
    # ×
    mx = ax + cell * N + Inches(0.22)
    _cap(s, mx, gy + cell * N / 2 - Inches(0.18), Inches(0.3), "×", 22, INK, PP_ALIGN.CENTER, True)
    # B
    bx = mx + Inches(0.45)
    for i in range(N):
        for j in range(N):
            f = PALE2 if j == 4 else WHITE
            rect(s, bx + j * cell, gy + i * cell, cell, cell, fill=f, line=LINE, lw=0.75)
    _cap(s, bx, gy - Inches(0.32), cell * N, "B  (K×N)", 13, ACCENT_D, PP_ALIGN.CENTER, True)
    _cap(s, bx, gy + cell * N + Inches(0.08), cell * N, "열 타일 하나를 통째로 읽는다", 10.5, MUTED, PP_ALIGN.CENTER)
    # =
    ex = bx + cell * N + Inches(0.22)
    _cap(s, ex, gy + cell * N / 2 - Inches(0.18), Inches(0.3), "=", 22, INK, PP_ALIGN.CENTER, True)
    # C
    cx = ex + Inches(0.45)
    for i in range(N):
        for j in range(N):
            f = RGBColor(0xC7, 0xE4, 0xCF) if (i == 2 and j == 4) else WHITE
            rect(s, cx + j * cell, gy + i * cell, cell, cell, fill=f, line=LINE, lw=0.75)
    _cap(s, cx, gy - Inches(0.32), cell * N, "C  (M×N)", 13, ACCENT_D, PP_ALIGN.CENTER, True)
    _cap(s, cx, gy + cell * N + Inches(0.08), cell * N, "타일 하나가 완성된다", 10.5, GREEN, PP_ALIGN.CENTER)

    ty = gy + cell * N + Inches(0.52)
    items = [
        "A의 한 행 타일을 온칩에 올려두면, B의 열 타일을 바꿔 가며 C의 한 행 전체를 만들 수 있다 → A를 N/T번 재사용한다.",
        "타일 크기를 T라 하면 HBM 읽기량이 O(M·N·K)에서 O(M·N·K / T)로 줄어든다. T가 클수록 좋지만 온칩 용량이 한계다.",
        "제약식: 타일 3개(A·B·C)를 동시에 담아야 하므로 3·T²·(원소 바이트) ≤ 온칩 용량.",
    ]
    d._bullets(s, [{"t": t} for t in items], ML, ty, BODY_W, BODY_BOT - ty - Inches(0.78))
    d._callout(s, "타일링은 '계산을 쪼개는 것'이 아니라 '데이터를 재사용하려고 계산 순서를 바꾸는 것'이다.")
    return s


# ---------------------------------------------------------------- 3. 메모리 계층
def memory_tiers(d):
    s = d._new()
    top = d._chrome(s, "메모리 계층 — 위로 갈수록 작고 빠르다", "RNGD 실제 수치")
    rows = [
        ("TRF  (Tensor Register File)", "8 KB / lane", "Contraction Engine의 정지 피연산자", 0.34, RGBColor(0xB9, 0xD5, 0xF5)),
        ("VRF  (Vector Register File)", "8 KB / slice", "Vector Engine이 매 사이클 읽는다", 0.42, RGBColor(0xC7, 0xDE, 0xF8)),
        ("DM  (Data Memory, 온칩 SRAM)", "512 KB / slice", "커널의 주 작업 메모리. 여기에 타일을 올린다", 0.62, RGBColor(0xD6, 0xE7, 0xFA)),
        ("SPM  (온칩 SRAM)", "임시·중간값", "컴파일러가 관리. API는 아직 제한적", 0.76, RGBColor(0xE4, 0xEF, 0xFC)),
        ("HBM  (패키지 위 대용량)", "47.5 GiB / 카드", "가중치·활성 장기 저장. 읽기 비용 최대", 1.0, RGBColor(0xF0, 0xF6, 0xFD)),
    ]
    y = top + Inches(0.28)
    maxw = BODY_W - Inches(4.9)
    for name, size, use, frac, fill in rows:
        w = maxw * frac
        x = ML + (maxw - w) / 2
        h = Inches(0.62)
        b = rect(s, x, y, w, h, fill=fill, line=LINE)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name
        style(r, 13, True, ACCENT_D)
        _cap(s, ML + maxw + Inches(0.22), y + Inches(0.06), Inches(1.25), size, 12, INK, PP_ALIGN.LEFT, True)
        dx = ML + maxw + Inches(1.52)
        _cap(s, dx, y + Inches(0.02), max(SW - MR - dx, Inches(1.0)), use, 10, MUTED)
        y += h + Inches(0.10)
    arrow(s, ML + Inches(0.12), y - Inches(0.1), ML + Inches(0.12), top + Inches(0.28), ACCENT, 2.0)
    _cap(s, ML - Inches(0.35), top + Inches(1.6), Inches(0.8), "빠름", 11, ACCENT, PP_ALIGN.CENTER, True)
    d._callout(s, "슬라이스 하나가 쓸 수 있는 온칩 메모리는 512 KB뿐이다. 4096×4096 fp32 행렬(64 MB)은 128조각으로 나눠야 들어간다.")
    return s


# ---------------------------------------------------------------- 4. 8단계 파이프라인
def pipeline(d):
    s = d._new()
    top = d._chrome(s, "텐서 유닛 8단계 파이프라인", "슬라이스마다 이 파이프라인이 하나씩 있다")
    stages = [
        ("Fetch", "DM에서\n스트림으로", PALE),
        ("Switch", "슬라이스 사이\n분배·전치", PALE),
        ("Collect", "32B flit로\n정규화", PALE),
        ("Contraction", "곱-합 축약\n(TRF 읽음)", PALE2),
        ("Vector", "원소 단위 연산\n(VRF 읽음)", PALE2),
        ("Cast", "dtype 변환", PALE3),
        ("Transpose", "flit 안 전치", PALE3),
        ("Commit", "DM에 쓰기", GREY),
    ]
    n = len(stages)
    gap = Inches(0.10)
    w = (BODY_W - gap * (n - 1)) / n
    y = top + Inches(0.75)
    h = Inches(1.55)
    for i, (nm, sub, fill) in enumerate(stages):
        x = ML + i * (w + gap)
        b = rect(s, x, y, w, h, fill=fill, line=LINE)
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.12)
        para(tf, True, nm, 12.5, True, ACCENT_D, align=PP_ALIGN.CENTER)
        for ln in sub.split("\n"):
            para(tf, False, ln, 9.5, False, INK, align=PP_ALIGN.CENTER, space_before=2)
        _cap(s, x, y - Inches(0.30), w, str(i + 1), 12, MUTED, PP_ALIGN.CENTER, True)
        if i < n - 1:
            arrow(s, x + w, y + h / 2, x + w + gap, y + h / 2, ACCENT, 1.4)
    # 메모리 표시
    ry = y + h + Inches(0.22)
    rect(s, ML, ry, w * 1.05, Inches(0.4), fill=WHITE, line=ACCENT)
    _cap(s, ML, ry + Inches(0.08), w * 1.05, "DM", 12, ACCENT_D, PP_ALIGN.CENTER, True)
    lastx = ML + (n - 1) * (w + gap)
    rect(s, lastx, ry, w, Inches(0.4), fill=WHITE, line=ACCENT)
    _cap(s, lastx, ry + Inches(0.08), w, "DM", 12, ACCENT_D, PP_ALIGN.CENTER, True)
    trx = ML + 3 * (w + gap)
    rect(s, trx, ry, w, Inches(0.4), fill=WHITE, line=AMBER)
    _cap(s, trx, ry + Inches(0.08), w, "TRF", 12, AMBER, PP_ALIGN.CENTER, True)
    vrx = ML + 4 * (w + gap)
    rect(s, vrx, ry, w, Inches(0.4), fill=WHITE, line=AMBER)
    _cap(s, vrx, ry + Inches(0.08), w, "VRF", 12, AMBER, PP_ALIGN.CENTER, True)

    ty = ry + Inches(0.62)
    items = [
        {"t": "순서가 고정돼 있다. 소프트웨어는 필요한 단계만 이어 붙이고 나머지는 건너뛴다. 순서를 바꿀 수는 없다."},
        {"t": "Switch만 슬라이스 경계를 넘는다. 나머지는 전부 자기 슬라이스 안에서만 움직인다."},
        {"t": "Collect가 어떤 크기의 패킷이든 32바이트 flit로 정규화한다. 뒤의 엔진들은 flit만 소비한다."},
    ]
    d._bullets(s, items, ML, ty, BODY_W, BODY_BOT - ty - Inches(0.78))
    d._callout(s, "이 파이프라인이 vISA 프로그래밍의 뼈대다. 커널을 쓴다는 것은 이 8단계를 어떻게 통과시킬지 적는 일이다.")
    return s


# ---------------------------------------------------------------- 5. 하드웨어 계층
def hierarchy(d):
    s = d._new()
    top = d._chrome(s, "하드웨어 계층 — Chip · Cluster · Slice · Lane", "우리 서버 기준 실측값")
    y = top + Inches(0.35)
    levels = [
        ("서버", "RNGD 4장 · 총 190 GiB HBM · 총 32 코어", RGBColor(0xEC, 0xF2, 0xF9), 1.0),
        ("Chip", "카드 1장 = 칩 1개 · HBM 47.5 GiB", RGBColor(0xDC, 0xE9, 0xF8), 0.84),
        ("Cluster", "칩당 2개", RGBColor(0xCB, 0xDF, 0xF6), 0.66),
        ("Slice", "클러스터당 256개 → 칩당 512개. 슬라이스 1개 = 텐서 유닛 1개 · DM 512 KB", RGBColor(0xB8, 0xD4, 0xF3), 0.48),
        ("Lane", "슬라이스당 최대 8개 → 칩당 4096개. MAC 배열의 한 행 · TRF 8 KB", RGBColor(0xA4, 0xC8, 0xF0), 0.30),
    ]
    ZONE = BODY_W * 0.58          # 막대는 왼쪽 58% 만 쓰고, 오른쪽은 설명 자리
    for nm, desc, fill, frac in levels:
        w = ZONE * frac
        x = ML + (ZONE - w) / 2
        h = Inches(0.66)
        b = rect(s, x, y, w, h, fill=fill, line=LINE)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = nm
        style(r, 15, True, ACCENT_D)
        _cap(s, ML, y + Inches(0.19), BODY_W, "", 10)
        dx = x + w + Inches(0.16)
        tf2 = tbox(s, dx, y + Inches(0.14), max(SW - MR - dx, Inches(0.6)), Inches(0.5))
        para(tf2, True, desc, 11, False, MUTED)
        y += h + Inches(0.12)
    d._callout(s, "Chip / Cluster / Slice / Lane 은 '공간 차원'이다. 매핑에서 이 네 축에 텐서를 어떻게 펼칠지 직접 적는다.")
    return s


# ---------------------------------------------------------------- 6. 매핑 Stride/Modulo
def mapping_grid(d):
    s = d._new()
    top = d._chrome(s, "★ Stride와 Modulo를 같이 쓰면 그것이 타일링이다",
                    "axes![A = 12] 을 4개짜리 블록 3개로 나눈다")
    cell = Inches(0.52)
    y = top + Inches(0.45)
    x0 = ML + Inches(0.5)
    # 원본 인덱스 줄
    _cap(s, ML, y - Inches(0.02), Inches(0.5), "A", 14, INK, PP_ALIGN.RIGHT, True)
    for i in range(12):
        b = rect(s, x0 + i * cell, y, cell, cell, fill=WHITE, line=LINE, lw=0.75)
        label(b, str(i), 12, False, INK)
    # A/4
    y2 = y + cell + Inches(0.30)
    _cap(s, ML, y2 - Inches(0.02), Inches(0.5), "A/4", 14, ACCENT_D, PP_ALIGN.RIGHT, True)
    cols = [PALE, PALE2, PALE3]
    for i in range(12):
        b = rect(s, x0 + i * cell, y2, cell, cell, fill=cols[i // 4], line=LINE, lw=0.75)
        label(b, str(i // 4), 12, True, ACCENT_D)
    _cap(s, x0 + 12 * cell + Inches(0.18), y2 + Inches(0.10), Inches(3.4),
         "블록 번호 (0,0,0,0, 1,1,1,1, 2,2,2,2)", 12, MUTED)
    # A%4
    y3 = y2 + cell + Inches(0.22)
    _cap(s, ML, y3 - Inches(0.02), Inches(0.5), "A%4", 14, AMBER, PP_ALIGN.RIGHT, True)
    for i in range(12):
        b = rect(s, x0 + i * cell, y3, cell, cell, fill=WHITE, line=LINE, lw=0.75)
        label(b, str(i % 4), 12, True, AMBER)
    _cap(s, x0 + 12 * cell + Inches(0.18), y3 + Inches(0.10), Inches(3.4),
         "블록 안 위치 (0,1,2,3 반복)", 12, MUTED)

    ty = y3 + cell + Inches(0.34)
    items = [
        {"t": "m![A/4] 은 '몇 번째 블록인가', m![A%4] 는 '블록 안 몇 번째인가'를 준다. 둘을 합치면 원래 인덱스가 복원된다.", "b": True},
        {"t": "항등식: E ≡ m![{E}/n, {E}%n]  — 어떤 축이든 스트라이드와 모듈로로 쪼갤 수 있고, 쪼개도 의미는 같다."},
        {"t": "쪼갠 두 조각을 서로 다른 하드웨어 차원에 배정하면 그게 타일링이다. 예: 블록 번호는 Time(순차 반복), 블록 안 위치는 Lane(동시 처리)."},
        {"t": "주의: m![A/n] 은 A의 크기가 n으로 나눠떨어져야 한다. 안 되면 컴파일 단계에서 막힌다. 그래서 Padding(#)이 필요해진다."},
    ]
    d._bullets(s, items, ML, ty, BODY_W, BODY_BOT - ty - Inches(0.78))
    d._callout(s, "4부에서 배운 타일링이 vISA에서는 새 문법이 아니라 '/'와 '%' 두 글자로 표현된다.")
    return s


# ---------------------------------------------------------------- 7. 이중 버퍼링
def double_buffer(d):
    s = d._new()
    top = d._chrome(s, "이중 버퍼링 — 가져오기와 계산을 겹친다", "TRF를 FirstHalf / SecondHalf로 나눠 번갈아 쓴다")
    y = top + Inches(0.55)
    unit = (BODY_W - Inches(1.5)) / 8
    x0 = ML + Inches(1.5)
    # 눈금
    for i in range(9):
        hline(s, x0 + i * unit, y - Inches(0.18), Emu(1), LINE)
        _cap(s, x0 + i * unit - Inches(0.2), y - Inches(0.42), Inches(0.4), f"t{i}", 9.5, MUTED, PP_ALIGN.CENTER)

    def track(ty, name, blocks, color):
        _cap(s, ML, ty + Inches(0.10), Inches(1.35), name, 12, INK, PP_ALIGN.RIGHT, True)
        for (st, ln, txt, fill) in blocks:
            b = rect(s, x0 + st * unit, ty, unit * ln, Inches(0.46), fill=fill, line=LINE)
            label(b, txt, 10.5, True, ACCENT_D)

    track(y, "순진한 방식", [
        (0, 2, "가져오기 0", PALE), (2, 2, "계산 0", PALE3),
        (4, 2, "가져오기 1", PALE), (6, 2, "계산 1", PALE3)], PALE)
    _cap(s, x0, y + Inches(0.52), BODY_W - Inches(1.5),
         "타일 2개 처리에 8단위. 가져오는 동안 연산기는 놀고, 계산하는 동안 DMA는 논다.", 11, RED)

    y2 = y + Inches(1.15)
    track(y2, "이중 버퍼링", [
        (0, 2, "가져오기 0 → Half A", PALE), (2, 2, "가져오기 1 → Half B", PALE),
        (4, 2, "가져오기 2 → Half A", PALE)], PALE)
    track(y2 + Inches(0.52), "", [
        (2, 2, "계산 0 (Half A)", PALE3), (4, 2, "계산 1 (Half B)", PALE3),
        (6, 2, "계산 2 (Half A)", PALE3)], PALE3)
    _cap(s, x0, y2 + Inches(1.04), BODY_W - Inches(1.5),
         "같은 일을 6단위에 끝낸다. 정상 상태에서는 DMA와 연산이 항상 동시에 돈다.", 11, GREEN)

    ty = y2 + Inches(1.40)
    items = [
        {"t": "핵심 조건: 지금 계산하는 버퍼와 지금 채우는 버퍼가 서로 달라야 한다. 같으면 해저드(WAR/RAW)가 나서 스케줄러가 기다리게 만든다."},
        {"t": "vISA에서는 TRF 주소를 FirstHalf / SecondHalf 로 지정해 이를 실현한다. 프리페치는 보통 sub 컨텍스트에 올린다."},
        {"t": "이득의 상한은 min(가져오기 시간, 계산 시간)이다. 우리 실측처럼 DMA가 훨씬 길면 겹쳐도 DMA 시간이 그대로 남는다."},
    ]
    d._bullets(s, items, ML, ty, BODY_W, BODY_BOT - ty - Inches(0.78))
    d._callout(s, "겹치기는 '느린 쪽을 빠르게' 만들지 못한다. 느린 쪽 뒤에 빠른 쪽을 숨길 뿐이다.")
    return s


# ---------------------------------------------------------------- 8. 추상화 사다리
def ladder(d):
    s = d._new()
    top = d._chrome(s, "vISA는 어느 높이의 언어인가", "위로 갈수록 편하고, 아래로 갈수록 통제권이 크다")
    rows = [
        ("PyTorch / 모델 코드", "무엇을 계산할지만 적는다", "레이아웃·스케줄 전부 위임", GREY),
        ("그래프 컴파일러 (furiosa-llm)", "연산 그래프를 자르고 붙인다", "커널 선택은 자동", GREY),
        ("TCL (커널 라이브러리 계층)", "제공되는 커널을 조합한다", "슬라이스 내부는 손댈 수 없다", PALE),
        ("★ vISA", "메모리 배치·스케줄·엔진 선택을 직접 적는다", "텐서 단위로 사고하되 통제권은 갖는다", PALE2),
        ("LIR → EDF (.bin)", "컴파일러 산출물", "사람이 쓰는 계층이 아니다", GREY),
        ("RNGD 하드웨어", "실제 실행", "", GREY),
    ]
    y = top + Inches(0.30)
    for i, (nm, what, note, fill) in enumerate(rows):
        h = Inches(0.64) if "vISA" in nm else Inches(0.54)
        b = rect(s, ML, y, BODY_W * 0.34, h, fill=fill,
                 line=ACCENT if "vISA" in nm else LINE, lw=2.0 if "vISA" in nm else 1.0)
        label(b, nm, 14 if "vISA" in nm else 12.5, True,
              ACCENT_D if "vISA" in nm else INK, align=PP_ALIGN.LEFT)
        tf = tbox(s, ML + BODY_W * 0.355, y + Inches(0.06), BODY_W * 0.33, h)
        para(tf, True, what, 11.5, False, INK)
        if note:
            tf = tbox(s, ML + BODY_W * 0.70, y + Inches(0.06), BODY_W * 0.30, h)
            para(tf, True, note, 11.5, False, MUTED)
        if i < len(rows) - 1:
            arrow(s, ML + BODY_W * 0.17, y + h, ML + BODY_W * 0.17, y + h + Inches(0.08),
                  RGBColor(0xA9, 0xB6, 0xC4), 1.2)
        y += h + Inches(0.08)
    d._callout(s, "vISA를 쓰는 이유는 하나다. TCL이 내려가지 못하는 슬라이스 내부까지 직접 지시하려는 것.")
    return s


# ---------------------------------------------------------------- 9. Roofline
def roofline(d):
    s = d._new()
    top = d._chrome(s, "Roofline — 내 커널은 어디에 있나", "가로축 연산강도, 세로축 도달 가능한 성능")
    ox, oy = ML + Inches(1.25), top + Inches(3.55)
    w, h = Inches(6.4), Inches(3.1)
    hline(s, ox, oy, w, INK, 1.5)
    from pptx.enum.shapes import MSO_CONNECTOR
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(ox), E(oy), E(ox), E(oy - h))
    c.line.color.rgb = INK
    c.line.width = Pt(1.5)
    knee_x = ox + w * 0.42
    top_y = oy - h * 0.78
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(ox), E(oy), E(knee_x), E(top_y))
    c.line.color.rgb = ACCENT
    c.line.width = Pt(3)
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(knee_x), E(top_y), E(ox + w), E(top_y))
    c.line.color.rgb = ACCENT
    c.line.width = Pt(3)
    _cap(s, ox - Inches(1.2), oy - h - Inches(0.02), Inches(1.15), "성능\n(FLOP/s)", 11, MUTED, PP_ALIGN.RIGHT)
    _cap(s, ox + w - Inches(2.2), oy + Inches(0.12), Inches(2.2), "연산강도 (FLOP/byte)", 11, MUTED, PP_ALIGN.RIGHT)
    _cap(s, ox + Inches(0.15), oy - h * 0.30, Inches(2.2), "메모리 바운드\n대역폭이 천장", 11, RED)
    _cap(s, knee_x + Inches(0.35), top_y - Inches(0.42), Inches(2.6), "연산 바운드\n연산기가 천장", 11, GREEN)
    for (fx, fy, nm, col) in ((0.14, 0.24, "GEMV\n(메모리 바운드)", RED),
                              (0.30, 0.52, "타일링 안 한 GEMM", AMBER),
                              (0.72, 0.78, "타일링한 GEMM", GREEN)):
        px, py = ox + w * fx, oy - h * fy
        dot = rect(s, px - Inches(0.07), py - Inches(0.07), Inches(0.14), Inches(0.14),
                   fill=col, shape=MSO_SHAPE.OVAL)
        _cap(s, px - Inches(1.0), py + Inches(0.10), Inches(2.0), nm, 10, col, PP_ALIGN.CENTER)
    bx = ox + w + Inches(0.55)
    items = [
        {"t": "연산강도 = 한 바이트를 가져와서 몇 번 연산하는가."},
        {"t": "왼쪽 기울어진 선은 메모리 대역폭이다. 여기 있으면 연산기를 아무리 늘려도 소용없다.", "lv": 1},
        {"t": "평평한 선은 연산기 최대 성능이다. 여기 닿아야 하드웨어를 다 쓴 것이다.", "lv": 1},
        {"t": "타일링의 목적은 점을 오른쪽으로 미는 것이다. 계산량은 그대로인데 이동량이 줄어드니 연산강도가 올라간다."},
        {"t": "우리 실측에서 vISA 커널 대부분은 왼쪽에 있다. 사이클의 96.5%가 DMA다(컴파일러 스케줄 모델 예측값)."},
    ]
    d._bullets(s, items, bx, top + Inches(0.35), SW - MR - bx, Inches(3.4))
    d._callout(s, "성능이 안 나올 때 물어야 할 첫 질문은 '연산기가 부족한가'가 아니라 '데이터를 몇 번 다시 읽고 있나'다.")
    return s


# ---------------------------------------------------------------- 10. 시스톨릭 어레이
def systolic(d):
    s = d._new()
    top = d._chrome(s, "MAC을 2차원으로 늘어놓으면 — 시스톨릭 어레이", "데이터가 흘러가며 부분합이 쌓인다")
    cell = Inches(0.72)
    N = 4
    gx, gy = ML + Inches(1.55), top + Inches(0.75)
    for i in range(N):
        for j in range(N):
            b = rect(s, gx + j * cell, gy + i * cell, cell - Inches(0.06), cell - Inches(0.06),
                     fill=PALE, line=ACCENT, lw=0.75)
            label(b, "MAC", 9.5, False, ACCENT_D)
    for i in range(N):
        arrow(s, gx - Inches(0.62), gy + i * cell + cell / 2 - Inches(0.03),
              gx - Inches(0.06), gy + i * cell + cell / 2 - Inches(0.03), AMBER, 1.4)
        _cap(s, gx - Inches(1.45), gy + i * cell + cell / 2 - Inches(0.18), Inches(0.8),
             f"A 행{i}", 10, AMBER, PP_ALIGN.RIGHT)
    for j in range(N):
        arrow(s, gx + j * cell + cell / 2 - Inches(0.03), gy - Inches(0.55),
              gx + j * cell + cell / 2 - Inches(0.03), gy - Inches(0.06), ACCENT, 1.4)
        _cap(s, gx + j * cell - Inches(0.1), gy - Inches(0.86), Inches(0.85),
             f"B 열{j}", 10, ACCENT, PP_ALIGN.CENTER)
    for j in range(N):
        arrow(s, gx + j * cell + cell / 2 - Inches(0.03), gy + N * cell - Inches(0.02),
              gx + j * cell + cell / 2 - Inches(0.03), gy + N * cell + Inches(0.45), GREEN, 1.4)
    _cap(s, gx, gy + N * cell + Inches(0.48), cell * N, "C 부분합이 아래로 빠져나온다", 10.5, GREEN, PP_ALIGN.CENTER)

    bx = gx + N * cell + Inches(0.85)
    items = [
        {"t": "MAC 하나는 a×b+c 만 한다. 곱하고 이전 값에 더한다."},
        {"t": "왼쪽에서 A의 행이, 위에서 B의 열이 흘러 들어온다. 각 MAC은 지나가는 값을 곱해 부분합에 더하고 옆·아래로 넘긴다."},
        {"t": "한 번 읽은 값이 배열 전체를 지나며 여러 번 쓰인다 → 읽기 한 번에 연산 여러 번. 연산강도를 하드웨어가 직접 올려준다.", "b": True},
        {"t": "대가: 모양이 배열 크기에 딱 맞아야 한다. 안 맞으면 패딩으로 채우고, 채운 칸은 버려야 한다."},
        {"t": "GPU의 Tensor Core도, TCP의 Contraction Engine도 이 아이디어의 변형이다."},
    ]
    d._bullets(s, items, bx, top + Inches(0.45), SW - MR - bx, Inches(4.3))
    d._callout(s, "'모양이 딱 맞아야 한다'는 제약이 타일링·패딩·정렬 문제의 뿌리다.")
    return s


# ---------------------------------------------------------------- 11. DMA 지배
def dma_share(d):
    s = d._new()
    top = d._chrome(s, "성능의 진실 — 사이클의 대부분은 데이터 이동이다",
                    "커널 130개의 컴파일러 스케줄을 합산한 값 (모델 예측, 벽시계 실측 아님)")
    y = top + Inches(0.55)
    bar_w = BODY_W - Inches(2.4)
    bx = ML + Inches(1.9)
    segs = [("DmaEngine  96.5%", 0.965, RGBColor(0xC1, 0x3C, 0x37)),
            ("PeCore 3.3%", 0.033, RGBColor(0x1A, 0x7F, 0x37)),
            ("", 0.002, GREY)]
    x = bx
    for nm, frac, col in segs:
        w = bar_w * frac
        b = rect(s, x, y, w, Inches(0.78), fill=col, line=WHITE)
        if nm:
            label(b, nm, 14 if frac > 0.1 else 9, True, WHITE)
        x += w
    _cap(s, ML, y + Inches(0.22), Inches(1.75), "전체 사이클", 13, INK, PP_ALIGN.RIGHT, True)

    y2 = y + Inches(1.25)
    rows = [("DmaEngine", "75,464,336 사이클", "명령 470개", "96.5%"),
            ("PeCore", "2,586,167 사이클", "명령 1,557개", "3.3%")]
    for i, (a, b_, c_, dd) in enumerate(rows):
        yy = y2 + i * Inches(0.46)
        rect(s, ML, yy, BODY_W, Inches(0.42), fill=WHITE if i % 2 == 0 else RGBColor(0xF6, 0xF9, 0xFC), line=LINE)
        for txt, xx, ww, col, bold in ((a, ML + Inches(0.16), Inches(2.2), INK, True),
                                       (b_, ML + Inches(2.5), Inches(2.6), INK, False),
                                       (c_, ML + Inches(5.3), Inches(2.2), MUTED, False),
                                       (dd, ML + Inches(7.8), Inches(1.4), ACCENT_D, True)):
            _cap(s, xx, yy + Inches(0.09), ww, txt, 12, col, PP_ALIGN.LEFT, bold)

    ty = y2 + Inches(1.05)
    items = [
        {"t": "명령 수는 PeCore가 3배 많은데(1,557 대 470) 사이클은 DMA가 29배 많다. 명령 하나당 걸리는 시간이 압도적으로 다르다."},
        {"t": "커널 하나하나로 봐도 같다. DMA 비중 중앙값 82.8%, 130개 중 107개가 50% 이상, 54개가 90% 이상이다."},
        {"t": "그래서 슬라이스 내부 연산을 아무리 잘 짜도 전체의 3.3% 안에서만 이득이 난다(암달의 법칙).", "b": True},
        {"t": "최적화 레버는 매핑·타일링·데이터 배치 쪽이다. 즉 4부와 7부에서 배운 것이 곧 성능 도구다."},
    ]
    d._bullets(s, items, ML, ty, BODY_W, BODY_BOT - ty - Inches(0.78))
    d._callout(s, "이 수치는 컴파일러 스케줄 모델의 예측이다. 벽시계로 검증한 값이 아니라는 점을 항상 같이 말해야 한다.")
    return s
