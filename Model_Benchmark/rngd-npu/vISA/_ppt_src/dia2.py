#!/usr/bin/env python3
"""추가 도해 — 글로만 설명하던 개념을 그림으로."""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from deck import (ACCENT, ACCENT_D, AMBER, BODY_BOT, BODY_W, CODE_BG, CODE_FG,
                  E, GREEN, INK, LINE, ML, MR, MUTED, RED, SW, SH, WHITE,
                  arrow, hline, label, para, rect, style, tbox)
from dia import GREY, PALE, PALE2, PALE3, _cap


def _grid(s, x, y, cell, rows, cols, fill=None, val=None, line=LINE, fs=10, fc=INK):
    """작은 격자. val(i,j) 로 칸 글자, fill(i,j) 로 칸 색."""
    for i in range(rows):
        for j in range(cols):
            f = fill(i, j) if callable(fill) else (fill or WHITE)
            b = rect(s, x + j * cell, y + i * cell, cell, cell, fill=f, line=line, lw=0.75)
            if val:
                v = val(i, j)
                if v:
                    label(b, str(v), fs, False, fc)


def _tail(d, s, items, y, callout=None):
    d._bullets(s, items, ML, y, BODY_W, BODY_BOT - y - (Inches(0.78) if callout else 0))
    if callout:
        d._callout(s, callout)


# ---------------------------------------------------------------- 1부
def contraction3(d):
    s = d._new()
    top = d._chrome(s, "텐서 축약은 늘 세 단계다 — 펼치고, 곱하고, 접는다",
                    "Broadcast → Multiply → Reduce")
    y = top + Inches(0.30)
    w = (BODY_W - Inches(1.2)) / 3
    steps = [
        ("① Broadcast (펼치기)", "없는 축을 복제해\n모양을 맞춘다", PALE,
         "A[i,k] → A[i,k,j]\nB[k,j] → B[i,k,j]"),
        ("② Multiply (곱하기)", "같은 자리끼리\n원소별로 곱한다", PALE2,
         "T[i,k,j] = A·B\n(i×k×j 개의 곱)"),
        ("③ Reduce (접기)", "축약할 축을 따라\n전부 더해 없앤다", PALE3,
         "C[i,j] = Σ_k T[i,k,j]\nk 축이 사라진다"),
    ]
    for n, (title, desc, fill, formula) in enumerate(steps):
        x = ML + n * (w + Inches(0.6))
        b = rect(s, x, y, w, Inches(1.05), fill=fill, line=LINE)
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.10)
        para(tf, True, title, 14, True, ACCENT_D, align=PP_ALIGN.CENTER)
        for ln in desc.split("\n"):
            para(tf, False, ln, 11, False, INK, align=PP_ALIGN.CENTER, space_before=2)
        fb = rect(s, x, y + Inches(1.15), w, Inches(0.72), fill=RGBColor(0xF7, 0xFA, 0xFD), line=LINE)
        tf = fb.text_frame
        for k, ln in enumerate(formula.split("\n")):
            para(tf, k == 0, ln, 11, k == 0, INK if k == 0 else MUTED,
                 mono=True, align=PP_ALIGN.CENTER, space_before=1)
        if n < 2:
            arrow(s, x + w + Inches(0.06), y + Inches(0.52), x + w + Inches(0.54), y + Inches(0.52))
    ty = y + Inches(2.05)
    _tail(d, s, [
        {"t": "행렬곱은 이 셋의 특수한 경우다. 축약할 축이 k 하나뿐인 경우다.", "b": True},
        {"t": "Attention 의 QKᵀ 도, Convolution 도, MLP 도 축약할 축만 다를 뿐 같은 세 단계다."},
        {"t": "TCP(Tensor Contraction Processor) 라는 이름이 여기서 나온다. 이 세 단계를 하드웨어로 직접 실행하려고 만든 칩이다."},
        {"t": "실제 하드웨어는 ①을 물리적으로 복제하지 않는다. 브로드캐스트 배선으로 같은 값을 여러 연산기에 동시에 흘려보낸다.", "lv": 1},
    ], ty, "축약 = 펼치고 · 곱하고 · 접기. 딥러닝 연산의 99% 가 이 틀에 들어간다.")
    return s


def gemm_ladder(d):
    s = d._new()
    top = d._chrome(s, "내적 → GEMV → GEMM — 같은 연산이 커지는 것뿐이다",
                    "재사용이 어디서 생기는지 보라")
    c = Inches(0.26)
    y = top + Inches(0.45)
    xs = [ML + Inches(0.4), ML + Inches(4.3), ML + Inches(8.2)]

    # 내적: 1×K · K×1 → 1
    x = xs[0]
    _grid(s, x, y, c, 1, 6, PALE)
    _cap(s, x, y - Inches(0.30), c * 6, "a  (1×K)", 11.5, ACCENT_D, PP_ALIGN.CENTER, True)
    _grid(s, x, y + Inches(0.55), c, 6, 1, PALE2)
    _cap(s, x + Inches(0.35), y + Inches(0.55) + c * 2, Inches(1.4), "· b (K×1)", 11.5, ACCENT_D)
    r1 = rect(s, x + Inches(1.5), y + Inches(1.0), c, c, fill=RGBColor(0xC7, 0xE4, 0xCF), line=LINE)
    _cap(s, x + Inches(1.35), y + Inches(1.32), Inches(0.9), "= 스칼라", 11, GREEN, PP_ALIGN.CENTER)
    _cap(s, x, y + Inches(2.35), Inches(3.4), "내적 (dot product)\n곱 K번 · 재사용 0", 12, INK, PP_ALIGN.LEFT, True)

    # GEMV: M×K · K×1 → M×1
    x = xs[1]
    _grid(s, x, y, c, 6, 6, PALE)
    _cap(s, x, y - Inches(0.30), c * 6, "A  (M×K)", 11.5, ACCENT_D, PP_ALIGN.CENTER, True)
    _grid(s, x + c * 6 + Inches(0.25), y, c, 6, 1, PALE2)
    _cap(s, x + c * 6 + Inches(0.15), y - Inches(0.30), Inches(0.6), "x", 11.5, ACCENT_D, PP_ALIGN.CENTER, True)
    _grid(s, x + c * 6 + Inches(0.85), y, c, 6, 1, RGBColor(0xC7, 0xE4, 0xCF))
    _cap(s, x + c * 6 + Inches(0.75), y - Inches(0.30), Inches(0.6), "y", 11.5, GREEN, PP_ALIGN.CENTER, True)
    _cap(s, x, y + Inches(2.35), Inches(3.4), "GEMV (행렬×벡터)\n곱 MK번 · x 를 M번 재사용", 12, INK, PP_ALIGN.LEFT, True)

    # GEMM: M×K · K×N → M×N
    x = xs[2]
    _grid(s, x, y, c, 6, 6, PALE)
    _cap(s, x, y - Inches(0.30), c * 6, "A  (M×K)", 11.5, ACCENT_D, PP_ALIGN.CENTER, True)
    _grid(s, x + c * 6 + Inches(0.25), y, c, 6, 6, PALE2)
    _cap(s, x + c * 6 + Inches(0.25), y - Inches(0.30), c * 6, "B  (K×N)", 11.5, ACCENT_D, PP_ALIGN.CENTER, True)
    _cap(s, x, y + Inches(2.35), Inches(4.0),
         "GEMM (행렬×행렬)\n곱 MNK번 · A 를 N번, B 를 M번 재사용", 12, INK, PP_ALIGN.LEFT, True)

    ty = y + Inches(3.0)
    _tail(d, s, [
        {"t": "셋 다 \"곱해서 더한다\"는 같은 연산이다. 다른 것은 피연산자의 축이 몇 개냐뿐이다."},
        {"t": "재사용 횟수가 결정적이다. 내적은 읽은 값을 한 번 쓰고 버리고, GEMM 은 A 의 한 원소를 N 번, B 의 한 원소를 M 번 쓴다.", "b": True},
        {"t": "그래서 GEMM 만 하드웨어를 꽉 채울 수 있다. GEMV 는 아무리 연산기를 늘려도 메모리가 못 따라온다(2부에서 숫자로 본다)."},
    ], ty)
    return s


def einsum_read(d):
    s = d._new()
    top = d._chrome(s, "einsum 표기 읽는 법 — 사라지는 축을 찾아라",
                    "IK, KJ → IJ  (이것이 GEMM 이다)")
    y = top + Inches(0.55)
    expr = [("I", ACCENT), ("K", RED), (",", MUTED), (" ", MUTED), ("K", RED), ("J", ACCENT),
            (" → ", MUTED), ("I", ACCENT), ("J", ACCENT)]
    x = ML + Inches(2.6)
    for ch, col in expr:
        w = Inches(0.52) if len(ch) <= 1 else Inches(1.30)   # 공백 포함 길이로 잡는다
        tf = tbox(s, x, y, w, Inches(0.7))
        para(tf, True, ch, 42, True, col, mono=True, align=PP_ALIGN.CENTER)
        x += w
    _cap(s, ML, y + Inches(0.86), BODY_W,
         "왼쪽 = 입력 텐서들의 축,  오른쪽 = 출력 텐서의 축", 13, MUTED, PP_ALIGN.CENTER)

    y2 = y + Inches(1.22)
    rows = [
        ("I", ACCENT, "양쪽 입력 중 하나에만 있고 출력에도 있다", "자유 축 — 그대로 남는다"),
        ("J", ACCENT, "양쪽 입력 중 하나에만 있고 출력에도 있다", "자유 축 — 그대로 남는다"),
        ("K", RED, "두 입력에 모두 있고 출력에는 없다", "★ 축약 축 — 이 축을 따라 더해서 없앤다"),
    ]
    for i, (nm, col, where, mean) in enumerate(rows):
        yy = y2 + i * Inches(0.52)
        b = rect(s, ML + Inches(0.6), yy, Inches(0.56), Inches(0.44), fill=WHITE, line=col, lw=2.0)
        label(b, nm, 20, True, col, mono=True)
        _cap(s, ML + Inches(1.45), yy + Inches(0.06), Inches(4.6), where, 12, INK)
        _cap(s, ML + Inches(6.2), yy + Inches(0.06), Inches(5.4), mean, 12, col, PP_ALIGN.LEFT, True)

    ty = y2 + Inches(1.70)
    _tail(d, s, [
        {"t": "규칙은 하나다. 출력에 없는 축이 축약 축이고, 그 축을 따라 합을 취한다."},
        {"t": "예: 배치 행렬곱 BIK, BKJ → BIJ 는 B 가 양쪽 입력과 출력에 다 있으니 축약하지 않고 그대로 반복한다."},
    ], ty, "축약 축이 몇 개인지, 그 축의 크기가 얼마인지가 곧 연산량이다.")
    return s


# ---------------------------------------------------------------- 2부
def mem_cost(d):
    s = d._new()
    top = d._chrome(s, "연산은 싸고 이동은 비싸다", "같은 값 하나를 다루는 데 드는 상대 비용 (자릿수 감각)")
    rows = [
        ("레지스터 / TRF·VRF 읽기", 1, RGBColor(0x1A, 0x7F, 0x37)),
        ("fp32 곱셈 1회", 2, RGBColor(0x1A, 0x7F, 0x37)),
        ("온칩 SRAM (DM) 읽기", 12, ACCENT),
        ("칩 안 다른 슬라이스로 이동", 30, AMBER),
        ("HBM 읽기", 200, RED),
        ("PCIe 로 호스트 왕복", 2000, RGBColor(0x7A, 0x10, 0x1A)),
    ]
    maxv = 2000
    y = top + Inches(0.20)
    barx = ML + Inches(3.6)
    barw = BODY_W - Inches(4.9)
    import math
    for nm, v, col in rows:
        _cap(s, ML, y + Inches(0.10), Inches(3.4), nm, 12, INK, PP_ALIGN.RIGHT, True)
        w = barw * (math.log10(v + 1) / math.log10(maxv + 1))
        rect(s, barx, y, max(w, Inches(0.12)), Inches(0.34), fill=col, line=None)
        _cap(s, barx + max(w, Inches(0.12)) + Inches(0.12), y + Inches(0.10), Inches(1.6),
             f"× {v:,}", 12, col, PP_ALIGN.LEFT, True)
        y += Inches(0.44)
    _cap(s, barx, y + Inches(0.00), barw, "가로축은 로그 눈금이다. 실제 차이는 그림보다 훨씬 크다.", 10, MUTED)
    ty = y + Inches(0.30)
    _tail(d, s, [
        {"t": "정확한 배수는 칩마다 다르지만 자릿수는 비슷하다. 연산 대 DRAM 접근이 대략 100배다.", "b": True},
        {"t": "타일링(4부)·이중 버퍼링·매핑 설계(7부)는 전부 이 표의 아래쪽 줄을 위쪽 줄로 바꾸려는 시도다."},
    ], ty, "성능을 물을 때는 \"연산이 몇 번인가\"가 아니라 \"HBM 을 몇 번 읽는가\"를 먼저 세라.")
    return s


def intensity(d):
    s = d._new()
    top = d._chrome(s, "연산강도 — 한 바이트를 가져와서 몇 번 계산하나",
                    "연산강도 = FLOP ÷ 이동 바이트")
    y = top + Inches(0.35)
    cases = [
        ("GEMV  (M=K=4096, fp32)", "2·M·K = 3,355만 FLOP",
         "(M·K + K + M)·4B ≈ 6,714만 B", "≈ 0.5 FLOP/B", RED, "메모리 바운드 — 연산기가 논다"),
        ("GEMM 타일링 없음 (M=N=K=4096)", "2·M·N·K = 1,374억 FLOP",
         "매 곱마다 A,B 를 새로 읽으면 ≈ 1,374억 B", "≈ 1 FLOP/B", AMBER, "여전히 메모리 바운드"),
        ("GEMM 타일 T=64", "2·M·N·K = 1,374억 FLOP",
         "이동량 ≈ 2·M·N·K/T · 4B", "≈ 32 FLOP/B", GREEN, "연산 바운드에 접근"),
    ]
    for nm, fl, by, ai, col, verdict in cases:
        rect(s, ML, y, BODY_W, Inches(0.88), fill=RGBColor(0xFA, 0xFC, 0xFE), line=LINE)
        _cap(s, ML + Inches(0.2), y + Inches(0.10), Inches(3.9), nm, 13, INK, PP_ALIGN.LEFT, True)
        _cap(s, ML + Inches(0.2), y + Inches(0.42), Inches(3.9), fl, 11, MUTED)
        _cap(s, ML + Inches(0.2), y + Inches(0.62), Inches(4.6), by, 10.5, MUTED)
        bb = rect(s, ML + Inches(5.2), y + Inches(0.16), Inches(2.3), Inches(0.56), fill=col, line=None)
        label(bb, ai, 17, True, WHITE)
        _cap(s, ML + Inches(7.8), y + Inches(0.30), Inches(4.2), verdict, 12.5, col, PP_ALIGN.LEFT, True)
        y += Inches(0.98)
    ty = y + Inches(0.02)
    _tail(d, s, [
        {"t": "같은 GEMM 인데 타일 크기 하나로 연산강도가 32배 달라진다. 계산량은 한 글자도 안 바뀌었다.", "b": True},
        {"t": "연산강도가 낮으면 아무리 좋은 연산기를 넣어도 소용없다. 데이터가 안 도착한다."},
    ], ty, "타일링은 연산을 줄이지 않는다. 같은 연산을 하면서 읽는 횟수만 줄인다.")
    return s


# ---------------------------------------------------------------- 3부
def mac_unit(d):
    s = d._new()
    top = d._chrome(s, "MAC — 모든 것의 기본 단위", "곱하고, 이전 값에 더한다. 그게 전부다.")
    cy = top + Inches(1.05)
    ax, bx = ML + Inches(1.2), ML + Inches(1.2)
    ai = rect(s, ax, cy - Inches(0.55), Inches(0.8), Inches(0.42), fill=PALE, line=ACCENT)
    label(ai, "a", 16, True, ACCENT_D, mono=True)
    bi = rect(s, bx, cy + Inches(0.25), Inches(0.8), Inches(0.42), fill=PALE, line=ACCENT)
    label(bi, "b", 16, True, ACCENT_D, mono=True)
    mul = rect(s, ML + Inches(2.7), cy - Inches(0.35), Inches(0.9), Inches(0.9),
               fill=PALE2, line=AMBER, shape=MSO_SHAPE.OVAL)
    label(mul, "×", 26, True, AMBER)
    arrow(s, ax + Inches(0.8), cy - Inches(0.34), ML + Inches(2.7), cy - Inches(0.10), ACCENT, 1.4)
    arrow(s, bx + Inches(0.8), cy + Inches(0.46), ML + Inches(2.7), cy + Inches(0.20), ACCENT, 1.4)
    add = rect(s, ML + Inches(4.6), cy - Inches(0.35), Inches(0.9), Inches(0.9),
               fill=PALE3, line=GREEN, shape=MSO_SHAPE.OVAL)
    label(add, "+", 26, True, GREEN)
    arrow(s, ML + Inches(3.6), cy + Inches(0.10), ML + Inches(4.6), cy + Inches(0.10), AMBER, 1.6)
    acc = rect(s, ML + Inches(4.75), cy - Inches(1.35), Inches(0.6), Inches(0.42), fill=GREY, line=MUTED)
    label(acc, "c", 15, True, INK, mono=True)
    arrow(s, ML + Inches(5.05), cy - Inches(0.93), ML + Inches(5.05), cy - Inches(0.37), MUTED, 1.4)
    out = rect(s, ML + Inches(6.3), cy - Inches(0.21), Inches(1.7), Inches(0.62),
               fill=RGBColor(0xC7, 0xE4, 0xCF), line=GREEN)
    label(out, "a×b + c", 15, True, GREEN, mono=True)
    arrow(s, ML + Inches(5.5), cy + Inches(0.10), ML + Inches(6.3), cy + Inches(0.10), GREEN, 1.6)
    _cap(s, ML + Inches(8.2), top + Inches(0.35), Inches(3.7),
         "누산기 c 가 핵심이다.\n곱셈 결과를 매번 밖으로 내보내지 않고\n안에 쌓아 두기 때문에\n메모리 왕복이 사라진다.", 12.5, INK)

    ty = cy + Inches(1.35)
    _tail(d, s, [
        {"t": "행렬곱의 안쪽 루프 C[i][j] += A[i][k] * B[k][j] 가 정확히 MAC 한 번이다.", "b": True},
        {"t": "누산은 보통 입력보다 넓은 타입으로 한다. bf16 두 개를 곱해 fp32 로 쌓는 식이다. 수천 번 더하면 오차가 쌓이기 때문이다."},
        {"t": "MAC 을 몇 개 넣느냐가 칩의 이론 성능을 정한다. RNGD 는 슬라이스마다 MAC 배열이 있고 레인 하나가 그 배열의 한 행이다."},
    ], ty, "\"몇 TFLOPS\" 는 결국 \"MAC 이 몇 개이고 몇 GHz 로 도나\" 다.")
    return s


def precision_bits(d):
    s = d._new()
    top = d._chrome(s, "정밀도 — 비트를 어디에 쓸 것인가", "지수(범위) 대 가수(정밀도)의 배분")
    y = top + Inches(0.40)
    unit = Inches(0.24)
    fmts = [
        ("fp32", 1, 8, 23, "표준. 학습에 쓴다"),
        ("tf32", 1, 8, 10, "NVIDIA. fp32 범위 + fp16 정밀도"),
        ("bf16", 1, 8, 7, "fp32 와 범위가 같다. 추론의 기본"),
        ("fp16", 1, 5, 10, "범위가 좁아 오버플로 주의"),
        ("fp8", 1, 4, 3, "최신 추론. 스케일링 필수"),
        ("int8", 0, 0, 8, "정수. 양자화 필요"),
    ]
    for nm, sgn, exp, man, note in fmts:
        _cap(s, ML, y + Inches(0.06), Inches(0.9), nm, 13, INK, PP_ALIGN.RIGHT, True)
        x = ML + Inches(1.15)
        for cnt, col, lab in ((sgn, RGBColor(0x4B, 0x5A, 0x6B), "S"),
                              (exp, AMBER, "지수"),
                              (man, ACCENT, "가수")):
            if cnt == 0:
                continue
            w = unit * cnt
            b = rect(s, x, y, w, Inches(0.34), fill=col, line=WHITE, lw=0.75)
            if w > Inches(0.5):
                label(b, f"{lab} {cnt}", 10.5, True, WHITE)
            x += w
        _cap(s, ML + Inches(1.15) + unit * 33 + Inches(0.2), y + Inches(0.05), Inches(3.4),
             note, 11, MUTED)
        y += Inches(0.46)
    ty = y + Inches(0.18)
    _tail(d, s, [
        {"t": "지수 비트가 표현 범위를, 가수 비트가 정밀도를 정한다. bf16 이 인기 있는 이유는 fp32 와 지수가 같아 범위 문제가 안 생기기 때문이다.", "b": True},
        {"t": "누산은 입력보다 넓게 한다. bf16 곱 → fp32 누산이 표준이다. 좁게 누산하면 수천 번 더할 때 값이 뭉개진다."},
        {"t": "vISA 에서 dtype 변환은 Cast Engine 이 담당하고, Fetch·Commit 단계에서도 변환할 수 있다."},
    ], ty)
    return s


# ---------------------------------------------------------------- 4부
def padding_lane(d):
    s = d._new()
    top = d._chrome(s, "나눠떨어지지 않으면 패딩이 생긴다", "그리고 패딩 칸은 반드시 계산에서 빼야 한다")
    c = Inches(0.40)
    y = top + Inches(0.45)
    _cap(s, ML, y - Inches(0.30), Inches(6.0), "원소 13개를 레인 8개 단위로 처리해야 한다면", 13, INK, PP_ALIGN.LEFT, True)
    for i in range(16):
        real = i < 13
        f = PALE if real else RGBColor(0xF3, 0xD9, 0xD9)
        b = rect(s, ML + (i % 8) * c, y + (i // 8) * (c + Inches(0.14)), c, c,
                 fill=f, line=LINE, lw=0.75)
        label(b, str(i) if real else "?", 11, not real, INK if real else RED)
    _cap(s, ML + c * 8 + Inches(0.3), y + Inches(0.05), Inches(4.6),
         "그룹 1 — 8개 전부 실제 데이터", 12, GREEN)
    _cap(s, ML + c * 8 + Inches(0.3), y + c + Inches(0.20), Inches(5.4),
         "그룹 2 — 5개만 진짜, 3개는 패딩(임의값)", 12, RED)

    y2 = y + 2 * (c + Inches(0.14)) + Inches(0.30)
    box = rect(s, ML, y2, BODY_W, Inches(0.92), fill=RGBColor(0xFF, 0xF7, 0xE6),
               line=RGBColor(0xF0, 0xD2, 0x9B))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.24)
    tf.margin_top = Inches(0.12)
    para(tf, True, "패딩 칸에는 임의의 값이 들어 있다. 0 이 아니다.", 14, True, AMBER)
    para(tf, False, "그대로 더하면 답이 틀린다. 그래서 하드웨어가 \"이 flit 에서 앞의 몇 개가 진짜인가\"를 세어 알려주는 장치가 필요하다 — 그것이 VCG(Valid Count Generator)다.",
         12, False, INK, space_before=4)

    ty = y2 + Inches(1.10)
    _tail(d, s, [
        {"t": "vISA 매핑에서 패딩은 `#` 연산자로 명시한다. m![A # 8] 은 \"A 를 8의 배수로 채운다\"는 뜻이다."},
        {"t": "실측: 부분 충전 그룹을 HBM 출력 타입에 그대로 노출하는 패턴은 컴파일 단계에서 거부된다. 실패값 228/232/240 이 정확히 7×32 + 4×레인 수와 맞는다."},
        {"t": "레인을 8개 미만으로 쓰는 것 자체는 금지가 아니다. MNIST 커널은 m![1] 로 실기에서 정상 동작한다.", "lv": 1},
    ], ty, "패딩은 공짜가 아니다. 채운 만큼 연산기가 헛돌고, 빼는 것을 잊으면 답이 틀린다.")
    return s


# ---------------------------------------------------------------- 5부
def flit_packet(d):
    s = d._new()
    top = d._chrome(s, "flit 과 packet — 데이터가 흐르는 단위", "Collect 가 무엇이 들어오든 32바이트로 정규화한다")
    y = top + Inches(0.50)
    bw = Inches(0.42)
    _cap(s, ML, y - Inches(0.32), Inches(5.0), "flit — 32 바이트", 14, ACCENT_D, PP_ALIGN.LEFT, True)
    for i in range(8):
        b = rect(s, ML + i * bw, y, bw, Inches(0.52), fill=PALE, line=LINE, lw=0.75)
        label(b, f"e{i}", 10.5, False, INK, mono=True)
    _cap(s, ML, y + Inches(0.58), bw * 8, "fp32 라면 원소 8개 (8 × 4B = 32B)", 10.5, MUTED, PP_ALIGN.CENTER)
    for i in range(16):
        b = rect(s, ML + Inches(4.2) + i * Inches(0.21), y, Inches(0.21), Inches(0.52),
                 fill=PALE3, line=LINE, lw=0.5)
    _cap(s, ML + Inches(4.2), y + Inches(0.58), Inches(3.36), "bf16 이라면 원소 16개 (16 × 2B = 32B)", 10.5, MUTED, PP_ALIGN.CENTER)

    y2 = y + Inches(1.25)
    _cap(s, ML, y2 - Inches(0.30), Inches(6.0), "packet — 64 바이트 = flit 2개", 14, ACCENT_D, PP_ALIGN.LEFT, True)
    for k in range(2):
        b = rect(s, ML + k * (bw * 8 + Inches(0.10)), y2, bw * 8, Inches(0.48),
                 fill=PALE2, line=AMBER)
        label(b, f"flit {k}", 12, True, AMBER)
    _cap(s, ML + Inches(7.5), y2 + Inches(0.10), Inches(4.5),
         "Contraction Engine 의 Outer 단계가 받는 단위", 12, MUTED)

    ty = y2 + Inches(0.85)
    _tail(d, s, [
        {"t": "Collect Engine 이 임의 크기의 패킷을 32바이트 경계에 맞춰 패딩하고 flit 경계에서 자른다. 그 뒤의 엔진은 전부 flit 만 다룬다.", "b": True},
        {"t": "그래서 dtype 이 바뀌면 flit 하나에 담기는 원소 수가 바뀐다. fp32 는 8개, bf16 은 16개, int8 은 32개다."},
        {"t": "매핑에서 Packet 차원의 크기를 잘못 잡으면 여기서 어긋난다. 실기 컴파일 실패 메시지에 'Collect time mismatch' 나 'packet mismatch' 가 자주 나오는 이유다."},
    ], ty, "32바이트 flit 이 이 칩의 데이터 낱알이다. 모든 모양 계산이 이 단위로 맞아떨어져야 한다.")
    return s


def compile_stages(d):
    s = d._new()
    top = d._chrome(s, "컴파일 단계 — 내 Rust 코드가 칩 명령이 되기까지", "어느 단계에서 막히는지가 오류 메시지 앞에 붙는다")
    stages = [
        ("Rust 소스", "#[device] 함수\nm![] 타입", GREY),
        ("MIR", "매핑·모양 검증\n엔진 배치", PALE),
        ("VISA", "가상 명령\n엔진별 연산", PALE2),
        ("LIR", "자원 할당\n버퍼 크기 확정", PALE3),
        ("EDF (.bin)", "실제 명령 스트림\npert-ipc", GREY),
    ]
    n = len(stages)
    gap = Inches(0.42)
    w = (BODY_W - gap * (n - 1)) / n
    y = top + Inches(0.55)
    h = Inches(1.25)
    for i, (nm, sub, fill) in enumerate(stages):
        x = ML + i * (w + gap)
        b = rect(s, x, y, w, h, fill=fill, line=LINE)
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.12)
        para(tf, True, nm, 14, True, ACCENT_D, align=PP_ALIGN.CENTER, mono=(i in (1, 2, 3)))
        for ln in sub.split("\n"):
            para(tf, False, ln, 10.5, False, INK, align=PP_ALIGN.CENTER, space_before=2)
        if i < n - 1:
            arrow(s, x + w + Inches(0.04), y + h / 2, x + w + gap - Inches(0.04), y + h / 2)
    errs = [(1, "mir: Collect time mismatch\nmir: commit_trim packet mismatch"),
            (3, "lir: incorrect buffer size\nstrides(...) is not aligned by 8")]
    for i, msg in errs:
        x = ML + i * (w + gap)
        eb = rect(s, x - Inches(0.15), y + h + Inches(0.22), w + Inches(0.3), Inches(0.66),
                  fill=RGBColor(0xFD, 0xEC, 0xEC), line=RGBColor(0xE8, 0xB6, 0xB6))
        tf = eb.text_frame
        for k, ln in enumerate(msg.split("\n")):
            para(tf, k == 0, ln, 8.5, False, RED, mono=True, align=PP_ALIGN.CENTER, space_before=1)
    ty = y + h + Inches(1.05)
    _tail(d, s, [
        {"t": "오류 메시지 맨 앞의 `mir:` / `lir:` 가 어느 단계에서 막혔는지 알려준다. 이걸 보면 원인 범위가 확 좁아진다.", "b": True},
        {"t": "`mir:` 단계 오류는 대개 매핑 모양이 안 맞는 것이고, `lir:` 단계 오류는 버퍼 크기·정렬 문제다."},
        {"t": "중요: 이 전체 과정은 `--backend npu` 에서만 돈다. typecheck·emulation 백엔드는 MIR 아래로 내려가지 않으므로 여기서 걸릴 오류를 못 잡는다."},
    ], ty, "컴파일이 통과해도 실기 실행이 보장되지는 않는다. 커널 로드에서 따로 죽는 경우가 있다(10부).")
    return s


# ---------------------------------------------------------------- 7부
def mapping_pair(d):
    s = d._new()
    top = d._chrome(s, "Pair `m![A, B]` — 두 축을 한 줄로 잇는다",
                    "왼쪽이 상위(느리게 변함), 오른쪽이 하위(빠르게 변함)")
    c = Inches(0.50)
    y = top + Inches(0.55)
    _cap(s, ML, y - Inches(0.32), Inches(6.0), "axes![A = 3, B = 4] 일 때  m![A, B] 는 크기 12", 13, INK, PP_ALIGN.LEFT, True)
    x0 = ML + Inches(0.9)
    for i in range(12):
        a, b = i // 4, i % 4
        col = [PALE, PALE2, PALE3][a]
        bx = rect(s, x0 + i * c, y, c, c, fill=col, line=LINE, lw=0.75)
        label(bx, str(i), 12, False, INK)
        bb = rect(s, x0 + i * c, y + c + Inches(0.08), c, c, fill=WHITE, line=LINE, lw=0.75)
        label(bb, f"{a},{b}", 11, True, ACCENT_D)
    _cap(s, ML, y + Inches(0.13), Inches(0.85), "버퍼 위치", 11, MUTED, PP_ALIGN.RIGHT, True)
    _cap(s, ML, y + c + Inches(0.21), Inches(0.85), "(A, B)", 11, ACCENT_D, PP_ALIGN.RIGHT, True)

    ty = y + 2 * c + Inches(0.35)
    _tail(d, s, [
        {"t": "규칙: 위치 i 에서 A = i / 4, B = i % 4 다. 오른쪽 축(B)이 먼저 다 돌고 나서 왼쪽 축(A)이 하나 넘어간다.", "b": True},
        {"t": "C 언어의 다차원 배열 arr[A][B] 와 같은 순서다. 마지막 축이 메모리에서 연속이다."},
        {"t": "Pair 는 우결합이다. m![A, B, C] 는 m![A, m![B, C]] 로 읽는다."},
        {"t": "항등원은 m![1] 이다. m![A, 1] 은 m![A] 와 같다. 크기 1짜리 축을 끼워 넣어도 아무것도 안 바뀐다.", "lv": 1},
        {"t": "이 순서를 뒤집고 싶으면 데이터를 옮기는 게 아니라 m![B, A] 라고 쓰면 된다. 그게 vISA 가 전치를 다루는 방식이다."},
    ], ty, "Pair 하나로 \"어느 축이 메모리에서 연속인가\"가 정해진다.")
    return s


def mapping_padding(d):
    s = d._new()
    top = d._chrome(s, "Padding `#n` 과 Resize `=n` — 늘리기와 줄이기",
                    "둘 다 크기를 바꾸지만 방향과 목적이 반대다")
    c = Inches(0.46)
    y = top + Inches(0.50)
    # Padding
    _cap(s, ML, y - Inches(0.30), Inches(5.6), "m![A # 8]   —  A=5 를 8로 채운다", 13, ACCENT_D, PP_ALIGN.LEFT, True)
    for i in range(8):
        real = i < 5
        b = rect(s, ML + Inches(0.2) + i * c, y, c, c,
                 fill=PALE if real else RGBColor(0xF3, 0xD9, 0xD9), line=LINE, lw=0.75)
        label(b, str(i) if real else "?", 11, not real, INK if real else RED)
    _cap(s, ML + Inches(0.2), y + c + Inches(0.06), c * 8,
         "뒤 3칸은 임의값. 하드웨어 단위에 맞추려고 채운 것", 10, MUTED, PP_ALIGN.CENTER)
    # Resize
    x2 = ML + Inches(6.3)
    _cap(s, x2, y - Inches(0.30), Inches(5.4), "m![A = 3]   —  A=5 를 3으로 줄인다", 13, ACCENT_D, PP_ALIGN.LEFT, True)
    for i in range(5):
        keep = i < 3
        b = rect(s, x2 + Inches(0.2) + i * c, y, c, c,
                 fill=PALE3 if keep else GREY, line=LINE, lw=0.75)
        label(b, str(i), 11, False, INK if keep else RGBColor(0xA6, 0xB0, 0xBA))
    _cap(s, x2 + Inches(0.2), y + c + Inches(0.06), c * 5 + Inches(1.1),
         "뒤 2칸은 논리적으로 없는 셈 친다(절단)", 10, MUTED, PP_ALIGN.LEFT)

    ty = y + c + Inches(0.55)
    _tail(d, s, [
        {"t": "Padding 은 하드웨어 요구 때문에 쓴다. 레인 8개, flit 32바이트처럼 딱 떨어져야 하는 단위에 맞추려면 채울 수밖에 없다.", "b": True},
        {"t": "채운 칸의 값은 정해져 있지 않다. 0 이라고 가정하면 안 된다. 리듀스에서 반드시 제외해야 한다."},
        {"t": "Resize 는 반대로 논리적 크기를 줄인다. 큰 버퍼의 앞부분만 쓰고 싶을 때 쓴다."},
        {"t": "실전 함정: m![A / 4] 는 A 가 4로 나눠떨어져야 한다. 안 되면 먼저 m![A # 8 / 4] 처럼 패딩을 걸어 나눠떨어지게 만든다."},
    ], ty)
    return s


# ---------------------------------------------------------------- 9부
def contexts_timeline(d):
    s = d._new()
    top = d._chrome(s, "실행 컨텍스트 4종 — 서로 병렬로 돈다",
                    "프로그래머가 어느 컨텍스트에 올릴지 정하고, 스케줄러가 시점을 정한다")
    y = top + Inches(0.55)
    unit = (BODY_W - Inches(1.7)) / 10
    x0 = ML + Inches(1.7)
    for i in range(11):
        _cap(s, x0 + i * unit - Inches(0.22), y - Inches(0.36), Inches(0.44), f"t{i}", 9.5, MUTED, PP_ALIGN.CENTER)
    tracks = [
        ("ctx.main", "주 연산 (Contraction·Vector)", [(2, 3, "GEMM 타일 0"), (5, 3, "GEMM 타일 1")], PALE2),
        ("ctx.sub", "프리페치 (다음 타일 준비)", [(0, 2, "TRF 채우기 0"), (3, 2, "TRF 채우기 1"), (6, 2, "채우기 2")], PALE),
        ("ctx.tdma", "텐서 DMA (HBM ↔ DM)", [(0, 3, "가중치 로드"), (4, 3, "다음 블록 로드")], PALE3),
        ("ctx.pdma", "PCIe DMA (호스트 ↔ 칩)", [(8, 2, "결과 회수")], GREY),
    ]
    for k, (nm, desc, blocks, col) in enumerate(tracks):
        yy = y + k * Inches(0.62)
        _cap(s, ML, yy + Inches(0.02), Inches(1.55), nm, 12, ACCENT_D, PP_ALIGN.RIGHT, True)
        _cap(s, ML, yy + Inches(0.27), Inches(1.55), desc, 7.5, MUTED, PP_ALIGN.RIGHT)
        rect(s, x0, yy, unit * 10, Inches(0.5), fill=RGBColor(0xFA, 0xFC, 0xFE), line=LINE, lw=0.5)
        for st, ln, txt in blocks:
            b = rect(s, x0 + st * unit, yy, unit * ln, Inches(0.5), fill=col, line=LINE)
            label(b, txt, 10, True, ACCENT_D)
    ty = y + 4 * Inches(0.62) + Inches(0.12)
    _tail(d, s, [
        {"t": "네 컨텍스트는 따로 돈다. main 이 계산하는 동안 sub 가 다음 타일을 TRF 에 채우고 tdma 가 HBM 에서 다음 블록을 끌어온다.", "b": True},
        {"t": "4부의 이중 버퍼링이 실제로는 main 과 sub 를 나눠 쓰는 것이다. 같은 주소를 건드리면 스케줄러가 자동으로 대기시킨다(해저드)."},
    ], ty, "컨텍스트를 안 나누면 전부 순차 실행된다. vISA 로 성능을 내는 첫걸음이 이 분리다.")
    return s


def hazards(d):
    s = d._new()
    top = d._chrome(s, "해저드 3종 — 스케줄러가 기다리는 이유",
                    "같은 주소를 두 연산이 건드리면 순서를 지켜야 한다")
    y = top + Inches(0.35)
    gap = Inches(0.30)
    w = (BODY_W - gap * 2) / 3
    cases = [
        ("RAW", "Read After Write", "쓰고 나서 읽는다",
         ["① X 에 결과를 쓴다", "② X 를 읽어 다음 계산"],
         "②는 ①이 끝날 때까지 기다려야 한다. 안 그러면 옛 값을 읽는다.", RED),
        ("WAR", "Write After Read", "읽는 중에 덮어쓴다",
         ["① X 를 읽는 중", "② X 에 새 값을 쓴다"],
         "②가 먼저 끝나면 ①이 새 값을 읽어버린다. 이중 버퍼링이 막는 게 바로 이것.", AMBER),
        ("WAW", "Write After Write", "연달아 쓴다",
         ["① X 에 쓴다", "② X 에 또 쓴다"],
         "순서가 뒤집히면 최종 값이 틀린다.", ACCENT),
    ]
    for i, (nm, en, ko, steps, why, col) in enumerate(cases):
        x = ML + i * (w + gap)
        hb = rect(s, x, y, w, Inches(0.52), fill=col, line=None)
        tf = hb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"{nm}   {en}"
        style(r, 14, True, WHITE)
        bodyb = rect(s, x, y + Inches(0.52), w, Inches(2.35), fill=RGBColor(0xFA, 0xFC, 0xFE), line=LINE)
        tf = bodyb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.14)
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.16)
        para(tf, True, ko, 13, True, INK)
        for stp in steps:
            para(tf, False, stp, 11.5, False, ACCENT_D, mono=False, space_before=6)
        para(tf, False, why, 11, False, MUTED, space_before=8)
    ty = y + Inches(3.05)
    _tail(d, s, [
        {"t": "vISA 는 이 셋을 프로그래머에게 떠넘기지 않는다. 스케줄러가 주소를 분석해 자동으로 대기를 넣는다.", "b": True},
        {"t": "대신 대가가 있다. 대기가 많으면 병렬성이 죽는다. 버퍼를 나눠 쓰면(FirstHalf/SecondHalf) 애초에 해저드가 안 생긴다."},
        {"t": "성능이 안 나올 때 --dump-schedule 로 스케줄을 떠서 어디서 기다리는지 보는 것이 표준 절차다."},
    ], ty)
    return s


# ---------------------------------------------------------------- 8부
def kernel_flow(d):
    s = d._new()
    top = d._chrome(s, "커널 하나가 도는 모습 — 텐서 타입이 바뀌어 간다",
                    "vISA 코드는 이 흐름을 그대로 적은 것이다")
    steps = [
        ("DmaTensor", "HBM 에 있는 원본", GREY),
        ("DmTensor", "DM 으로 내려온 타일", PALE),
        ("StreamTensor", "Fetch 가 만든 스트림", PALE),
        ("CollectTensor", "32B flit 로 정규화", PALE),
        ("TrfTensor", "TRF 에 올린 가중치", PALE2),
        ("ContractTensor", "축약 결과(누산값)", PALE2),
        ("VectorTensor", "활성함수 적용", PALE3),
        ("DmTensor", "Commit 이 DM 에 씀", GREY),
    ]
    n = len(steps)
    y = top + Inches(0.16)
    h = Inches(0.42)
    for i, (t, desc, fill) in enumerate(steps):
        yy = y + i * (h + Inches(0.06))
        b = rect(s, ML + Inches(0.3), yy, Inches(3.1), h, fill=fill, line=LINE)
        label(b, t, 12.5, True, ACCENT_D, mono=True)
        _cap(s, ML + Inches(3.6), yy + Inches(0.09), Inches(3.4), desc, 11, INK)
        if i < n - 1:
            arrow(s, ML + Inches(1.85), yy + h, ML + Inches(1.85), yy + h + Inches(0.055),
                  ACCENT, 1.3)
    ops = [(1, "to_dm()"), (2, "fetch()"), (3, "collect()"), (4, "to_trf()"),
           (5, "contract()"), (6, "vector_*()"), (7, "commit()")]
    for i, nm in ops:
        yy = y + i * (h + Inches(0.06))
        ob = rect(s, ML + Inches(7.2), yy + Inches(0.03), Inches(1.9), Inches(0.36),
                  fill=WHITE, line=ACCENT)
        label(ob, nm, 10.5, True, ACCENT, mono=True)
        arrow(s, ML + Inches(7.15), yy + Inches(0.21), ML + Inches(6.9), yy + Inches(0.21),
              ACCENT, 1.1)
    _cap(s, ML + Inches(9.35), y + Inches(0.4), Inches(2.6),
         "타입이 바뀔 때마다\n메서드를 한 번 부른다.\n\n타입 시스템이 순서를\n강제하므로, 잘못된\n순서는 컴파일이 안 된다.", 12, INK)
    d._callout(s, "vISA 커널을 읽는 법: 타입 이름만 따라가면 데이터가 어느 메모리에 있는지 항상 알 수 있다.")
    return s


# ---------------------------------------------------------------- 10부
def results_matrix(d):
    s = d._new()
    top = d._chrome(s, "실측 결과 한눈에", "우리 서버 RNGD 4장 · 2026-07 · 프로세스 격리 실행")
    y = top + Inches(0.30)
    # 컴파일 매트릭스
    _cap(s, ML, y, Inches(6.0), "① 커널 200개를 npu 백엔드로 컴파일", 14, INK, PP_ALIGN.LEFT, True)
    bw = BODY_W - Inches(2.2)
    yy = y + Inches(0.34)
    segs = [("성공 137 (68.5%)", 137 / 200, GREEN), ("실패 63", 63 / 200, RED)]
    x = ML
    for nm, fr, col in segs:
        b = rect(s, x, yy, bw * fr, Inches(0.54), fill=col, line=WHITE)
        label(b, nm, 12.5, True, WHITE)
        x += bw * fr
    sub = [("진짜 로워링 공백", 24, RGBColor(0xC1, 0x3C, 0x37)),
           ("의도적 실패 표본", 23, RGBColor(0xD9, 0x7B, 0x4A)),
           ("컴파일러 ICE", 13, RGBColor(0xA8, 0x4B, 0x8C)),
           ("기타", 4, MUTED)]
    yy2 = yy + Inches(0.62)
    x = ML + bw * (137 / 200)
    for nm, v, col in sub:
        w = bw * (v / 200)
        b = rect(s, x, yy2, w, Inches(0.34), fill=col, line=WHITE, lw=0.5)
        x += w
    _cap(s, ML + bw * (137 / 200) - Inches(1.2), yy2 + Inches(0.38), Inches(5.6),
         "실패 63 = 로워링 공백 24 · 의도적 표본 23 · 컴파일러 ICE 13 · 기타 4", 10.5, MUTED)

    # 실기 매트릭스
    y3 = yy2 + Inches(0.86)
    _cap(s, ML, y3, Inches(6.0), "② 실기 테스트 89개를 각각 별도 프로세스로 실행", 14, INK, PP_ALIGN.LEFT, True)
    yy3 = y3 + Inches(0.34)
    segs2 = [("PASS 80", 80 / 89, GREEN), ("FAIL 5", 5 / 89, RED),
             ("ABORT 3", 3 / 89, RGBColor(0x7A, 0x10, 0x1A)), ("", 1 / 89, GREY)]
    x = ML
    for nm, fr, col in segs2:
        b = rect(s, x, yy3, bw * fr, Inches(0.54), fill=col, line=WHITE)
        if nm:
            label(b, nm, 12.5 if fr > 0.1 else 9, True, WHITE)
        x += bw * fr
    _cap(s, ML + bw + Inches(0.2), yy3 + Inches(0.14), Inches(1.9), "정상 동작 83\n(93.3%)", 12, GREEN, PP_ALIGN.LEFT, True)

    # 격리 효과
    y4 = yy3 + Inches(0.80)
    _cap(s, ML, y4, Inches(8.0), "③ 프로세스 격리 여부가 결과를 바꾼다 (vector_engine 테스트)", 14, INK, PP_ALIGN.LEFT, True)
    yy4 = y4 + Inches(0.32)
    for k, (lab, npass, col) in enumerate((("한 프로세스", 10, RED), ("테스트마다 격리", 33, GREEN))):
        b = rect(s, ML + Inches(2.4), yy4 + k * Inches(0.44), bw * (npass / 36) * 0.9, Inches(0.34),
                 fill=col, line=None)
        _cap(s, ML, yy4 + k * Inches(0.44) + Inches(0.05), Inches(2.25), lab, 12, INK, PP_ALIGN.RIGHT, True)
        _cap(s, ML + Inches(2.5) + bw * (npass / 36) * 0.9, yy4 + k * Inches(0.44) + Inches(0.05),
             Inches(2.0), f"{npass} 통과", 12, col, PP_ALIGN.LEFT, True)
    d._callout(s, "hang 커널 하나가 HAL -110 으로 같은 프로세스의 뒤 커널을 전부 오염시킨다. 측정 방법이 결론을 3배 바꾼다.")
    return s
