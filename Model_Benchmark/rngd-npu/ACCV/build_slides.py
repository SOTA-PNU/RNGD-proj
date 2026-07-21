#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""논문 그림·표 전부를 슬라이드로.
- 방법/어텐션 그림 = 이미지(다이어그램·히트맵이라 네이티브 차트 불가)
- 수치 결과 = 네이티브 편집가능 PPT 차트(파워포인트에서 데이터 수정 가능)
- 표 = 네이티브 PPT 표
번호는 논문 등장순서와 일치(그림1 방법, 그림2 어텐션, 그림3 처리량 / 표1~7).
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from PIL import Image

SW, SH = 13.333, 7.5
INK=RGBColor(0x22,0x22,0x22); SUB=RGBColor(0x66,0x66,0x66); BLUE=RGBColor(0x00,0x72,0xB2)
HEADBG=RGBColor(0x2b,0x3a,0x55); HEADTX=RGBColor(0xff,0xff,0xff); ALT=RGBColor(0xf2,0xf5,0xf8)

# 논문 identity 색
COL={"ToMe":RGBColor(0xD5,0x5E,0x00),"PiToMe":RGBColor(0x00,0x9E,0x73),
     "Ours":RGBColor(0x00,0x72,0xB2),"PiToMe+reg":RGBColor(0x00,0x72,0xB2),
     "무작위":RGBColor(0x99,0x99,0x99),"에너지":RGBColor(0x00,0x9E,0x73),
     "고노름":RGBColor(0xCC,0x79,0xA7),"no-reg":RGBColor(0x99,0x99,0x99)}
def scol(name):
    for k,v in COL.items():
        if name==k: return v
    return RGBColor(0x55,0x55,0x55)

prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH); blank=prs.slide_layouts[6]

def tb(slide,x,y,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Malgun Gothic"

def head(slide,title,sub):
    tb(slide,0.7,0.4,SW-1.4,0.85,title,28,INK,bold=True)
    if sub: tb(slide,0.7,1.18,SW-1.4,0.9,sub,13,SUB)

def fig_slide(title,sub,img):
    s=prs.slides.add_slide(blank); head(s,title,sub)
    iw,ih=Image.open(img).size; ar=iw/ih; mw,mh=SW-1.4,SH-2.3; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    s.shapes.add_picture(img,Inches((SW-w)/2),Inches(2.05+(mh-h)/2),Inches(w),Inches(h))

def chart_slide(title,sub,cats,series,xlab,ylab,ymin=None,ymax=None,
                ctype=XL_CHART_TYPE.LINE_MARKERS):
    """series = [(name,[vals]), ...] 네이티브 편집가능 차트."""
    s=prs.slides.add_slide(blank); head(s,title,sub)
    cd=CategoryChartData(); cd.categories=cats
    for name,vals in series: cd.add_series(name,vals)
    gf=s.shapes.add_chart(ctype,Inches(1.6),Inches(2.05),Inches(SW-3.2),Inches(SH-2.55),cd)
    ch=gf.chart
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False
    ch.legend.font.size=Pt(13); ch.legend.font.name="Malgun Gothic"
    for ser in ch.series:
        c=scol(ser.name)
        if ctype in (XL_CHART_TYPE.LINE_MARKERS,XL_CHART_TYPE.LINE):
            ser.format.line.color.rgb=c; ser.format.line.width=Pt(2.5); ser.smooth=False
            try:
                ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb=c
                ser.marker.format.line.color.rgb=RGBColor(0xff,0xff,0xff)
            except Exception: pass
        else:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb=c
    va=ch.value_axis; ca=ch.category_axis
    if ymin is not None: va.minimum_scale=ymin
    if ymax is not None: va.maximum_scale=ymax
    va.has_major_gridlines=True; va.has_title=True; va.axis_title.text_frame.text=ylab
    ca.has_title=True; ca.axis_title.text_frame.text=xlab
    for ax in (va,ca):
        ax.tick_labels.font.size=Pt(12); ax.tick_labels.font.name="Malgun Gothic"
        ax.axis_title.text_frame.paragraphs[0].runs[0].font.size=Pt(13)

# ---------- 표 파싱 ----------
def clean(c):
    c=re.sub(r"~?\\cite\{[^}]*\}","",c); c=re.sub(r"\\textbf\{([^}]*)\}",r"\1",c); c=re.sub(r"\\emph\{([^}]*)\}",r"\1",c)
    c=c.replace("${\\sim}","~").replace("{\\sim}","~").replace("${=}","=").replace("$+$","+").replace("$-$","-")
    c=c.replace("\\%","%").replace("\\,","").replace("$\\downarrow$","↓").replace("{>}",">").replace("{<}","<")
    c=re.sub(r"\$([^$]*)\$",r"\1",c); c=c.replace("{,}",",").replace("\\","").replace("textbf","").replace("underline","")
    c=c.replace("{","").replace("}","").replace("`","").replace("$","")
    return c.strip()

def parse_tables(md):
    out={}
    for blk in re.findall(r"\\begin\{table\}.*?\\end\{table\}",md,flags=re.S):
        cap=re.search(r"\\caption\{(.*?)\}\s*\\label",blk,flags=re.S)
        lab=re.search(r"\\label\{(tab:[^}]+)\}",blk)
        tab=re.search(r"\\begin\{tabular\}\{[^}]*\}(.*?)\\end\{tabular\}",blk,flags=re.S)
        if not(lab and tab): continue
        rows=[]
        for line in tab.group(1).split("\\\\"):
            line=re.sub(r"\\(top|mid|bottom)rule","",line).strip()
            if "&" not in line: continue
            rows.append([clean(x) for x in line.split("&")])
        out[lab.group(1)]=(clean(cap.group(1)) if cap else "", rows)
    return out

def col(rows,header,name):
    j=header.index(name)
    vals=[]
    for r in rows:
        try: vals.append(float(r[j]))
        except (ValueError,IndexError): vals.append(None)
    return vals

def short_cap(c):
    out=""
    for part in c.split("."):
        if out and len(out)+len(part)>150: break
        out+=part+"."
    return (out.strip() or c[:150]).rstrip(".")+"."

def table_slide(title,cap,rows):
    s=prs.slides.add_slide(blank); head(s,title,short_cap(cap))
    nrow=len(rows); ncol=max(len(r) for r in rows)
    rows=[r+[""]*(ncol-len(r)) for r in rows]
    tw=min(SW-1.4,2.0*ncol); th=min(SH-2.6,0.42*nrow)
    gt=s.shapes.add_table(nrow,ncol,Inches((SW-tw)/2),Inches(2.15),Inches(tw),Inches(th)).table
    for j in range(ncol): gt.columns[j].width=Inches(tw/ncol)
    for i,r in enumerate(rows):
        for j in range(ncol):
            cell=gt.cell(i,j); cell.text=r[j]
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
            run=p.runs[0] if p.runs else p.add_run()
            run.font.size=Pt(12); run.font.name="Malgun Gothic"
            if i==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=HEADBG; run.font.color.rgb=HEADTX; run.font.bold=True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb=(ALT if i%2 else RGBColor(0xff,0xff,0xff)); run.font.color.rgb=INK

# ================= 빌드 =================
md=open("main.md",encoding="utf-8").read()
T=parse_tables(md)

# 표지
s=prs.slides.add_slide(blank)
tb(s,0.9,2.5,SW-1.8,1.3,"레지스터 인지 토큰 축소",40,INK,bold=True,align=PP_ALIGN.CENTER)
tb(s,0.9,3.8,SW-1.8,0.7,"논문 그림·표 전부 · ACCV 2026",20,BLUE,align=PP_ALIGN.CENTER)

# ---- 그림(이미지): 방법, 어텐션 ----
fig_slide("그림 1 · 방법","각 블록에서 레지스터·CLS를 보호하고 패치만 병합. 깊이에 따라 ToMe는 레지스터를 잃고 Ours는 보존한다.","fig_method.png")
fig_slide("그림 2 · 어텐션 몰림","클래스 토큰의 어텐션 몰림을 색블록으로(빨강=많이·파랑=적게). Ours는 무압축과 비슷한 구조를 유지.","fig_attention.png")

# ---- 결과 그래프(네이티브 편집가능 차트) ----
# 주 결과: ToMe / PiToMe / Ours 정확도 vs 축소율
cats=[r[0] for r in T["tab:main"][1][1:]]
h_m=T["tab:main"][1][0]; h_p=T["tab:pitome"][1][0]
chart_slide("결과 그래프 1 · 주 결과","kNN top-1 정확도 vs 토큰 축소율. 축소가 커질수록 Ours의 이득이 벌어진다.",
    cats,[("ToMe",col(T["tab:main"][1][1:],h_m,"ToMe")),
          ("PiToMe",col(T["tab:pitome"][1][1:],h_p,"PiToMe")),
          ("Ours",col(T["tab:main"][1][1:],h_m,"Ours"))],
    "토큰 축소율","kNN top-1 (%)",ymin=68,ymax=82)

# 보호 기준 ablation: 같은 개수에서 무작위/에너지/고노름 vs 레지스터(Ours)
h_a=T["tab:ablation"][1][0]; ra=T["tab:ablation"][1][1:]
chart_slide("결과 그래프 2 · 보호 기준 ablation","같은 개수를 보호할 때 무작위·에너지·고노름은 ToMe와 비슷하고 레지스터(Ours)만 크게 오른다.",
    [r[0] for r in ra],[("ToMe",col(ra,h_a,"ToMe")),("무작위",col(ra,h_a,"무작위")),
        ("에너지",col(ra,h_a,"에너지")),("고노름",col(ra,h_a,"고노름")),("Ours",col(ra,h_a,"Ours"))],
    "토큰 축소율","kNN top-1 (%)",ymin=68,ymax=82)

# PiToMe+reg: 레지스터 보호가 PiToMe 병합 위에서도 이득
h_g=T["tab:generality"][1][0]; rg=T["tab:generality"][1][1:]
chart_slide("결과 그래프 3 · PiToMe+reg","레지스터 보호를 PiToMe 병합 위에 얹어도 모든 축소율에서 정확도가 오른다.",
    [r[0] for r in rg],[("PiToMe",col(rg,h_g,"PiToMe")),("PiToMe+reg",col(rg,h_g,"PiToMe+reg"))],
    "토큰 축소율","kNN top-1 (%)",ymin=70,ymax=82)

# (처리량은 표 tab:throughput 로 이동 — 선이 겹쳐 표가 읽기 좋음)

# 밀집 예측 mIoU
h_d=T["tab:dense"][1][0]; rd=T["tab:dense"][1][1:]
chart_slide("결과 그래프 4 · 밀집 예측 (ADE20k)","선형 세그멘테이션 mIoU. 레지스터 보호(Ours)가 극한 축소에서도 mIoU를 지킨다.",
    [r[0] for r in rd],[("ToMe",col(rd,h_d,"ToMe")),("무작위",col(rd,h_d,"무작위")),
        ("에너지",col(rd,h_d,"에너지")),("고노름",col(rd,h_d,"고노름")),("Ours",col(rd,h_d,"Ours"))],
    "토큰 축소율","mIoU (%)",ymin=18,ymax=30)

# val-LOO 불변성
h_c=T["tab:consistency"][1][0]; rc=T["tab:consistency"][1][1:]
chart_slide("결과 그래프 5 · val-LOO 불변성","갤러리를 val로 바꾼 leave-one-out에서도 순서와 이득이 유지된다.",
    [r[0] for r in rc],[("ToMe",col(rc,h_c,"ToMe")),("PiToMe",col(rc,h_c,"PiToMe")),
        ("Ours",col(rc,h_c,"Ours"))],
    "토큰 축소율","kNN top-1 (%)",ymin=60,ymax=78)

# DINOv3·ViT-5 확장: 극한(~95%)에서 레지스터 제거 시 붕괴 (막대). no-reg 있는 모델만.
re_=T["tab:extra"][1][1:]; cur=None; ex={}
for r in re_:
    if len(r)<2: continue
    if r[1]=="Ours": cur=r[0]; ex[cur]={"ours":float(r[-1]),"noreg":None}
    elif r[1]=="no-reg" and cur:
        try: ex[cur]["noreg"]=float(r[-1])
        except ValueError: pass
bar_m=[m for m in ex if ex[m]["noreg"] is not None]
chart_slide("결과 그래프 6 · DINOv3·ViT-5 확장","극한 축소(~95%)에서 레지스터를 제거(no-reg)하면 DINOv3는 거의 붕괴, Ours는 정확도 유지.",
    bar_m,[("Ours",[ex[m]["ours"] for m in bar_m]),("no-reg",[ex[m]["noreg"] for m in bar_m])],
    "모델","kNN top-1 (%)",ymin=0,ymax=90,ctype=XL_CHART_TYPE.COLUMN_CLUSTERED)

# ---- 표 전부(논문 번호) ----
titles={"tab:main":"표 1 · 주 결과 (train-갤러리 kNN)","tab:pitome":"표 2 · PiToMe 같은 예산",
"tab:ablation":"표 3 · 보호 기준 ablation","tab:throughput":"표 4 · GPU 처리량 (im/s)",
"tab:generality":"표 5 · 다른 병합기 (PiToMe+reg)","tab:extra":"표 6 · DINOv3·ViT-5 확장",
"tab:dense":"표 7 · 밀집 예측 (ADE20k)","tab:consistency":"표 8 · val-LOO 불변성"}
for lab in ["tab:main","tab:pitome","tab:ablation","tab:throughput","tab:generality","tab:extra","tab:dense","tab:consistency"]:
    if lab in T:
        cap,rows=T[lab]; table_slide(titles[lab],cap,rows)

prs.save("그림_슬라이드.pptx")
print("saved 그림_슬라이드.pptx · 슬라이드", len(prs.slides._sldIdLst))
