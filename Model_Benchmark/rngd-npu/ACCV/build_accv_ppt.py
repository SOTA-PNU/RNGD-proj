#!/usr/bin/env python3
"""ACCV 발표용 PPT 생성 — 네이티브 편집가능 차트/표(이미지 아님) + 해석.
python-pptx의 CategoryChartData를 쓰므로 데이터가 pptx 내장 워크시트에 들어가,
PowerPoint에서 수치를 직접 바꾸면 그래프가 갱신된다. 표도 네이티브(편집가능).
수치는 감사 통과한 보수적 값(단일 seed·단일 모델·kNN)만. 출력: ACCV_발표_결과.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUT = "/home/jun/RNGD-proj/Model_Benchmark/rngd-npu/ACCV/ACCV_발표_결과.pptx"
NAVY = RGBColor(0x1F, 0x2D, 0x3D); RED = RGBColor(0xD6, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66); BLUE = RGBColor(0x1F, 0x77, 0xB4)

# ---- 보수적 검증 데이터 (원자료: register_token_reduction/results, ablation/results) ----
RED_PCT = ['37%', '55%', '74%', '83%', '~91%']           # 실현 압축률(마지막 ~91%)
TOME = [74.70, 72.93, 70.34, 67.62, 63.91]
OURS = [75.71, 75.30, 74.13, 73.28, 71.98]
FULL = 76.35
ABL = {  # ~91% 압축, 동일 개수 keep-prior (단일 seed)
    'ToMe(무보호)': 63.90, 'Ours(register)': 71.98, 'Random': 63.19,
    'Energy(프록시)': 64.04, 'High-norm': 63.84}
ABL_ROWS = [  # 전 압축률 5전략
    ('37%', 74.70, 75.71, 74.70, 74.58, 74.78),
    ('55%', 72.95, 75.30, 72.82, 72.70, 72.87),
    ('74%', 70.36, 74.14, 70.00, 70.16, 70.16),
    ('83%', 67.60, 73.25, 67.26, 67.51, 67.53),
    ('~91%', 63.90, 71.98, 63.19, 64.04, 63.84)]
NPU = [('0(full)', 7.52), ('37%', 8.61), ('55%', 8.36), ('74%', 8.36), ('92%', 7.54)]  # 전구간 no speedup; r18(83%)만 컴파일FAIL
# reg 개수 스윕(k=0=ToMe … k=4=Ours) + 부트스트랩 95% CI (전체 50k GPU)
REG_SWEEP = [  # (압축, k0, k1, k2, k3, k4, CI_lo, CI_hi)  전체 50k
    ('55%', 72.93, 73.09, 75.13, 75.27, 75.13, 1.95, 2.46),
    ('74%', 70.14, 70.34, 74.20, 74.11, 74.21, 3.77, 4.37),
    ('91%', 64.01, 64.01, 72.15, 71.85, 71.89, 7.53, 8.25)]
# 검색 mAP(두 번째 지표, 전체 50k) — kNN 외 랭킹 지표서도 register 보호 우세
RETR_PCT = ['37%', '55%', '74%', '83%', '91%']
RETR_TOME = [51.46, 48.13, 43.39, 39.67, 34.57]
RETR_OURS = [54.23, 53.21, 51.45, 50.07, 47.69]

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8)).text_frame
    p = tb.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    return s


def bullets(s, items, left=0.6, top=1.15, width=6.0, height=5.6, size=15):
    tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.level = lvl
        p.font.size = Pt(size - lvl); p.font.color.rgb = NAVY if lvl == 0 else GRAY
        p.space_after = Pt(4)
    return tf


def add_line_chart(s, cats, series, left, top, w, h, title):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series: cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(left), Inches(top), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = True; ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10)
    return ch


def add_bar_chart(s, cats, name, vals, left, top, w, h, title):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(top), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = True; ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].font.size = Pt(12); ch.has_legend = False
    return ch


def add_table(s, rows, left, top, w, h, header=True, fs=11):
    nr, nc = len(rows), len(rows[0])
    tbl = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(w), Inches(h)).table
    for r in range(nr):
        for c in range(nc):
            cell = tbl.cell(r, c); cell.text = str(rows[r][c])
            para = cell.text_frame.paragraphs[0]; para.font.size = Pt(fs)
            if header and r == 0:
                para.font.bold = True; para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    return tbl


# ===== 1. 타이틀 =====
s = slide("Register-Aware Token Reduction — 결과·해석")
bullets(s, [
    ("register(전역정보) 토큰을 보호하며 중복 patch를 병합 → 극단 압축서 정확도 보존 (재학습 없음)", 0),
    ("본 슬라이드: 지금까지의 검증된(보수적) 결과 + 왜 이렇게 구성했고 무엇을 시사하는가", 1),
    ("⚠️ 모든 수치 = 단일 모델(DINOv2-ViT-B/14-reg4) · 단일 seed · kNN(k=20) — 일반화는 진행중", 1),
    ("표·그래프는 전부 PPT 내장 편집가능(수치 바꾸면 그래프 갱신)", 1),
], top=2.2, size=17)

# ===== 2. 배경·왜 이 실험인가 =====
s = slide("① 배경 — 왜 이 실험을 구성했나")
bullets(s, [
    ("문제: DINOv2 같은 인코더는 토큰 수백 개 처리 → 느림. 토큰 축소(병합)로 가속.", 0),
    ("기존 방법(ToMe=유사도, PiToMe=에너지)은 '무엇을 지킬지'를 일반 신호로만 정함.", 0),
    ("관찰: 일부 인코더는 전역정보를 소수 'register 토큰'에 몰아둠(Darcet ICLR'24).", 0),
    ("가설: 표준 병합은 register를 '배경처럼' 보고 합쳐 없앤다 → 극단 압축서 정확도 급락.", 0),
    ("→ 그래서 '같은 병합에서 register만 보호'를 통제 비교로 실험 설계(원인 격리).", 0),
    ("무엇을 쟀나: ImageNet kNN 정확도를 압축률별로, 보호대상만 바꿔 측정.", 0),
], top=1.3, width=12.0, size=17)

# ===== 3. 주 결과 =====
s = slide("② 주 결과 — register 보호가 극단 압축서 정확도 보존")
add_line_chart(s, RED_PCT, [('ToMe(무보호)', TOME), ('Ours(register보호)', OURS)],
               6.6, 1.2, 6.4, 4.4, "kNN top-1 (%) vs 토큰 축소율  (full=76.35)")
add_table(s, [['축소율', 'ToMe', 'Ours', 'Δ']] +
          [[RED_PCT[i], f"{TOME[i]:.2f}", f"{OURS[i]:.2f}", f"+{OURS[i]-TOME[i]:.2f}"] for i in range(5)],
          6.9, 5.7, 5.8, 1.3, fs=10)
bullets(s, [
    ("무엇을: DINOv2-reg, 같은 size-가중 ToMe식 병합, 보호대상만 CLS↔CLS+register.", 0),
    ("결과: 모든 압축률서 Ours 우세, 격차가 압축 강할수록 커짐(+1.0→+8.1).", 0),
    ("~91% 축소서 Ours 71.98(full −4.4) vs 무보호 63.91(full −12.4).", 0),
    ("시사: 전역정보 토큰(register)을 지키면 훨씬 더 줄여도 버팀.", 0),
    ("보수 표기: 단일 seed. +8p는 노이즈(±0.05)보다 2자릿수 큼(확고).", 1),
    ("압축률=명목 스케줄; 실현 ~91%, Ours가 register 4개만큼 토큰 약간 더 남김.", 1),
], top=1.25, width=6.2, size=14)

# ===== 4. Ablation =====
s = slide("③ 왜 'register라서'인가 — 동일 개수 keep-prior 비교")
add_bar_chart(s, list(ABL.keys()), 'kNN@~91%', list(ABL.values()), 6.7, 1.2, 6.3, 3.5,
              "~91% 축소서 보호전략별 kNN top-1 (%)")
add_table(s, [['축소율', 'ToMe', 'Ours', 'Rand', 'Energy', 'HiNorm']] +
          [[r[0]] + [f"{v:.1f}" for v in r[1:]] for r in ABL_ROWS], 6.7, 5.0, 6.3, 1.7, fs=9)
bullets(s, [
    ("무엇을: 같은 병합·같은 개수 보호, '무엇을 보호'만 변경(register/무작위/에너지/고노름).", 0),
    ("결과: register만 크게 이김(~91%서 71.98). 나머지는 무보호(63.9)와 노이즈(±0.9) 내.", 0),
    ("시사: 이득의 원인이 '토큰 더 보호'가 아니라 register 그 자체.", 0),
    ("⚠️ 정직 caveat(꼭 반영):", 0),
    ("Energy는 PiToMe식 프록시일 뿐, 실제 PiToMe 방법과 head-to-head 아님.", 1),
    ("★위 표는 '입력단 고정' 선택. 매 블록 동적 재선택하면 highnorm/energy 강해짐", 1),
    ("  (전체 50k: 극단서 ours 71.9 vs 동적 67.4 vs 정적/무보호 ~64) — ours +4.4 확실.", 1),
    ("★실제 PiToMe(공식) head-to-head, 전체 50k: 온건압축선 PiToMe 우세,", 0),
    ("  극단선 ours 우위(92%서 +3.2, 교차점~70%). 헤드라인은 '극단압축 우위'로 정직히 좁힘.", 1),
    ("★정식 ToMe(proportional attn+key유사도+attn↔MLP병합) head-to-head, 전체 50k:", 0),
    ("  ours 전압축률 이김 +2.2/+4.8/+10.3(55/74/91%) → '베이스라인이 약해서' 우려도 해소.", 1),
], top=1.25, width=6.3, size=13)

# ===== 4.7 reg 개수 스윕 + 부트스트랩 CI (인과·유의성) =====
s = slide("③-b 인과·유의성 — register 개수 스윕 + 부트스트랩 CI")
add_bar_chart(s, ['k=0(ToMe)', 'k=1', 'k=2', 'k=3', 'k=4(Ours)'], 'kNN@91%',
              list(REG_SWEEP[2][1:6]), 7.0, 1.25, 5.8, 3.3, "~91% 압축서 보호 register 개수 k별 kNN(%)")
add_table(s, [['압축', 'k=0', 'k=1', 'k=2', 'k=3', 'k=4', '95%CI(ours-tome)']] +
          [[r[0]] + [f"{v:.1f}" for v in r[1:6]] + [f"[{r[6]:+.1f},{r[7]:+.1f}]"] for r in REG_SWEEP],
          6.9, 5.1, 6.0, 1.2, fs=9)
bullets(s, [
    ("무엇을: 같은 병합, 보호하는 register 개수만 k=0→4로 늘림(k=0=ToMe, k=4=Ours).", 0),
    ("왜: '이득이 register 자체냐, 그냥 토큰 더 남겨서냐'를 가르고, 단일 seed 유의성을 확인.", 0),
    ("결과1(인과): k=1은 거의 무효, k≥2서 계단식 상승 후 포화 — '토큰 수'가 아니라", 0),
    ("  register를 묶음으로 지켜야 정보가 삶. 원인=register(정직: 엄밀 단조는 아님).", 1),
    ("결과2(유의성): 평가셋 부트스트랩 95% CI. 극단서 +7.88 [+7.53,+8.25]로 0을 넘김", 0),
    ("  전체 50k선 중간압축(55%)도 CI 0 배제[+2.0,+2.5] → 전 압축률 유의(이득 크기는 극단서 최대).", 1),
    ("시사: '단일 seed·토큰수 confound' 두 반론을 데이터로 방어.", 0),
    ("⚠️ 전체 50k 실측(GPU).", 1),
], top=1.3, width=6.6, size=13)

# ===== 3.5 검색 mAP (두 번째 지표) =====
s = slide("②-b 두 번째 지표 — 검색 mAP도 register 보호 우세")
add_line_chart(s, RETR_PCT, [('ToMe(무보호)', RETR_TOME), ('Ours(register보호)', RETR_OURS)],
               6.7, 1.25, 6.2, 4.4, "이미지 검색 mAP(%) vs 토큰 축소율 (전체 50k)")
add_table(s, [['축소율', 'ToMe', 'Ours', 'Δ']] +
          [[RETR_PCT[i], f"{RETR_TOME[i]:.1f}", f"{RETR_OURS[i]:.1f}", f"+{RETR_OURS[i]-RETR_TOME[i]:.1f}"] for i in range(5)],
          6.9, 5.75, 5.8, 1.2, fs=10)
bullets(s, [
    ("무엇을: 같은 특징으로 kNN이 아닌 표준 검색 지표(mAP) 측정.", 0),
    ("  각 query를 gallery 전체에 유사도 랭킹, 같은 클래스면 정답 → average precision 평균.", 1),
    ("왜: 지금까지 지표가 kNN 하나뿐 → '지표 하나 아니냐' 반론 대비(두 번째 축).", 0),
    ("결과: Ours 전 압축률 우세, 격차 +2.8→+13.1(압축 강할수록 커짐).", 0),
    ("  극단서 격차 +13.1 > kNN의 +7.9 → 우위가 kNN 특정 지표의 산물 아님.", 1),
    ("시사: register 보호가 만드는 건 '특정 지표 점수'가 아니라 더 나은 특징 그 자체.", 0),
    ("⚠️ 전체 50k 실측(GPU).", 1),
], top=1.3, width=6.4, size=13)

# ===== 4.5 메커니즘·효율 (GPU 불요, 직접 증거) =====
s = slide("④ 메커니즘·효율 — 직접 증거 (GPU 불요, 로컬)")
add_bar_chart(s, ['37%', '55%', '74%', '83%', '~91%'], 'FLOP 비율(full=1)',
              [0.825, 0.738, 0.653, 0.611, 0.569], 7.0, 1.3, 5.8, 3.2, "토큰 축소율별 백본 FLOP 비율")
add_table(s, [['축소율', 'FLOP절감', 'ours vs tome'],
              ['74%', '34.7%', '동일'], ['~91%', '43.1%', '동일(<0.1%)']], 7.2, 4.9, 5.4, 1.0, fs=10)
bullets(s, [
    ("메커니즘(병합빈도 직접 측정, r=16, 16장):", 0),
    ("표준 ToMe서 register의 94%가 병합돼 사라짐(평균 3.6/12블록, 이르게).", 1),
    ("Ours는 보호라 100% 생존 → '병합이 register를 없앤다'를 간접추론 아닌 직접 증거로.", 1),
    ("효율(FLOP, 계산):", 0),
    ("토큰 축소가 FLOP 17~43% 절감(91%축소서 43%). 91%축소≠91%절감(MLP 지배로 완만).", 1),
    ("★Ours·ToMe FLOP 사실상 동일(<0.1%) → +8%는 '토큰 더 남겨서'가 아님(compute 공짜).", 1),
    ("단 NPU 실측 지연은 안 줄음(별개) — 효율은 FLOP에 한정, 정직.", 1),
    ("왜 이 실험: '병합이 register 파괴'(메커니즘)+'이득이 공짜'(효율)를 데이터로 못박기.", 0),
], top=1.3, width=6.8, size=12.5)

# ===== 5. 대조군 + NPU (정직) =====
s = slide("④ 대조군·NPU 지연 — 정직한 한계")
add_bar_chart(s, [n for n, _ in NPU], 'NPU ms/forward', [v for _, v in NPU], 7.0, 1.3, 5.8, 3.2,
              "RNGD 실칩 지연 (낮을수록 빠름) — 축소해도 안 줄어듦")
bullets(s, [
    ("음성 대조군(register 없는 DINOv2): Ours==ToMe가 구성상 동일 → Δ≈0.", 0),
    ("  = 노이즈 바닥·하네스 검증일 뿐, '원인=register' 확증은 ③ ablation 몫.", 1),
    ("  주의: 이 모델은 원래 압축에 강건(무보호도 −4p) → register 외 차이도 있음.", 1),
    ("NPU 실칩 지연(rngd:0): 92% 축소서도 7.54ms=full(7.52) 그대로 — 전 구간 속도이득 0(0.87~1.0x).", 0),
    ("  (r=18/83%만 컴파일 실패, furiosa 미구현; r=20/92%은 컴파일되고도 이득 0.)", 1),
    ("  원인 추정: 이 토큰수 구간은 고정 오버헤드 지배 + 정적 슬라이스 측정.", 1),
    ("  → 논문서 'NPU 가속' 주장 안 함(정직). 효율은 FLOP/토큰 감소로 별도 제시.", 1),
], top=1.3, width=6.7, size=13)

# ===== 5.5 dense (분할) — GPU 결과 =====
s = slide("⑤ dense(분할) mIoU — register가 가장 빛나는 곳")
add_line_chart(s, ['37%', '55%', '74%', '83%', '~91%'],
               [('ToMe(무보호)', [21.12, 19.31, 17.32, 15.71, 14.05]),
                ('Ours(register)', [22.42, 21.71, 20.43, 19.12, 17.19])],
               6.6, 1.2, 6.4, 4.2, "ADE20k mIoU(%) vs 토큰 축소율 (full=23.1)")
add_table(s, [['축소율', 'ToMe', 'Ours', 'Rand', 'Energy', 'HiNorm'],
              ['74%', '17.3', '20.4', '16.9', '16.9', '16.5'],
              ['~91%', '14.1', '17.2', '13.4', '13.1', '13.3']], 6.7, 5.6, 6.3, 1.1, fs=9)
bullets(s, [
    ("무엇을: DINOv2-reg에 선형 seg head, 토큰 병합·unmerge 후 ADE20k mIoU. 5전략 비교.", 0),
    ("결과: Ours가 전 압축률 1위, 격차 압축 강할수록↑(~91%서 +3.1).", 0),
    ("★분류와 달리 dense선 random/energy/highnorm이 전혀 안 도움(무보호 이하) — register만 효과.", 0),
    ("시사: 공간정보 중요한 dense일수록 register 보호 이득 크고 깨끗함(가설 확증).", 0),
    ("정직: mIoU 절대값 낮음(16×16 저해상 probing)이나 상대비교 유효. n=2000·단일 seed.", 1),
], top=1.3, width=6.2, size=13)

# ===== 6. 한계·향후 =====
s = slide("⑥ 한계와 향후 실험 (보수적 로드맵)")
bullets(s, [
    ("현재 근거의 한계(과장 금지):", 0),
    ("모델 의존적(전체 50k·3모델): base +7.9 뚜렷, large는 무보호 붕괴(2.8)로 상대이득 큼, small은 동적energy가 근소우위(61.5 vs 60.7).", 1),
    ("주비교 베이스라인 ToMe-style(양팔 동일). 결정적이라 seed 대신 부트스트랩. dense는 저해상 probing.", 1),
    ("완료된 보강(전부 전체 50k):", 0),
    ("★실제 PiToMe: 온건압축 PiToMe 우세·극단 ours(+3.2, 교차점~70%) + 처리량 ours 약우위.", 1),
    ("★정식 ToMe(proportional attn 포함)로도 ours 전압축률 승(극단 +10.3) — '약한 베이스라인' 해소.", 1),
    ("★검색 mAP(극단 +13.1>kNN) · 부트스트랩 CI 전압축률 0배제(55%도) · 3모델·동적ablation·linear-probe.", 1),
    ("dense mIoU(ours 1위, 대안 무효) · 병합빈도(register 94% 소멸) · FLOP(43%,동일).", 1),
    ("남은 것(선택, GPU): native 해상도(448/518) · 대규모 throughput. 핵심 실험은 전부 전체 50k 완료.", 0),
], top=1.2, width=12.2, size=15)

prs.save(OUT)
print("SAVED", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
