#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""fig_method(직관 버전) 다이어그램을 PowerPoint '네이티브 객체'로 재구성 (PNG 통삽입 아님).
모든 토큰/박스/링/화살표/라벨이 각각 편집 가능한 도형·텍스트박스·연결선이라
PowerPoint에서 직접 위치·색·글자를 바꿀 수 있다.
좌표는 fig_method.py 의 matplotlib 좌표계(x∈[0,130], y∈[0,60], y-up)를 슬라이드 inch(y-down)로 사상.
구성 = 범례 / (좌)One block: Attn->Merge->MLP + protect-pass·merge / (우)Across blocks: ToMe(레지스터 소멸) vs Ours(그룹링 유지).
출력: ACCV_method_편집용.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

OUT = "/home/jun/RNGD-proj/Model_Benchmark/rngd-npu/ACCV/ACCV_method_편집용.pptx"

# ---- 색 (fig_method.py 와 동일) ----
C_CLS  = RGBColor(0x2b, 0x3a, 0x55)
C_REG  = RGBColor(0x00, 0x72, 0xB2)
C_PATCH= RGBColor(0xdc, 0xdc, 0xda)
C_PROT = RGBColor(0xE6, 0x9F, 0x00)
INK    = RGBColor(0x22, 0x22, 0x22)
MUT    = RGBColor(0x8a, 0x8a, 0x8a)
ARROW  = RGBColor(0x5b, 0x6b, 0x7a)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BOXF   = RGBColor(0xf3, 0xf3, 0xf3)
BOXFM  = RGBColor(0xe7, 0xef, 0xf7)
BOXE   = RGBColor(0xbc, 0xbc, 0xbc)
CALL   = RGBColor(0x9d, 0xb4, 0xc9)
DIVID  = RGBColor(0xd2, 0xd2, 0xd2)
TOME_C = RGBColor(0x44, 0x44, 0x44)
GRAY66 = RGBColor(0x66, 0x66, 0x66)
REGBLU = RGBColor(0x0f, 0x5f, 0x9c)

# ---- 좌표 사상: mpl(x∈0..130, y∈0..60, y-up) -> inch(y-down) ----
S     = 0.096
X_OFF = 0.45
YMAX  = 60.0
Y_TOP = 1.35

def ix(x): return X_OFF + x * S
def iytop(yv): return Y_TOP + (YMAX - yv) * S

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6]); shp = s.shapes

def rrect(x, y, w, h, fill=None, line=None, line_w=1.0, round_amt=0.2):
    sp = shp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ix(x)), Inches(iytop(y + h)),
                       Inches(w * S), Inches(h * S))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    try: sp.adjustments[0] = round_amt
    except Exception: pass
    return sp

def text_at(x, y, text, size, color=INK, bold=False, italic=False, ha='left', va='center', w=None):
    lh = (size / 72.0) * 1.35
    # 글자 길이에 맞춘 '딱 맞는' 박스 폭(inch). 고정 8인치 박스 방지 → 선택 박스가 글자 크기와 같아짐.
    factor = 0.66 if bold else 0.6
    tw = w if w is not None else max(len(text) * size * factor / 72.0 + 0.04, 0.10)
    yc = iytop(y)
    top = yc - lh / 2
    left = ix(x) if ha == 'left' else (ix(x) - tw / 2 if ha == 'center' else ix(x) - tw)
    tb = shp.add_textbox(Inches(left), Inches(top), Inches(tw), Inches(lh))
    tf = tb.text_frame; tf.word_wrap = False; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[ha]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = 'Arial'
    return tb

def line(x1, y1, x2, y2, color=ARROW, width=1.6, dash=None, arrow=False):
    c = shp.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(ix(x1)), Inches(iytop(y1)),
                          Inches(ix(x2)), Inches(iytop(y2)))
    c.line.color.rgb = color; c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dash: ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    if arrow: ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c

W, H, GAP = 2.3, 2.7, 0.55

def tok(x, y, w=W, h=H, fc=C_PATCH, ec=None, label='', tc=WHITE, fs=7):
    rrect(x, y, w, h, fill=fc, line=ec, line_w=1.0, round_amt=0.22)
    if label:
        text_at(x + w / 2, y + h / 2, label, fs, color=tc, bold=True, ha='center', va='center')

def group_ring(x0, y, n, w=W, gap=GAP, h=H):
    span = n * (w + gap) - gap
    rrect(x0 - 0.7, y - 0.7, span + 1.4, h + 1.4, fill=None, line=C_PROT, line_w=2.1, round_amt=0.18)

def row(x0, y, specs, prot_n=0):
    x = x0
    for kind in specs:
        if kind == 'cls':   tok(x, y, fc=C_CLS, label='CLS', fs=5.6)
        elif kind == 'reg': tok(x, y, fc=C_REG, label='R', fs=7)
        else:               tok(x, y, fc=C_PATCH)
        x += W + GAP
    if prot_n > 0: group_ring(x0, y, prot_n)
    return x

# ================= 범례 =================
ly = 55.5
def legend_item(x, kind, txt):
    if kind == 'cls':   tok(x, ly, fc=C_CLS, label='CLS', fs=5.6)
    elif kind == 'reg': tok(x, ly, fc=C_REG, label='R', fs=7)
    else:               tok(x, ly, fc=C_PATCH)
    text_at(x + 3.1, ly + 1.35, txt, 10, INK, ha='left')
legend_item(38, 'cls', 'class token')
legend_item(60, 'reg', 'register')
legend_item(82, 'patch', 'patch')
tok(101, ly, fc=C_REG, label='R', fs=7); group_ring(101, ly, 1)
text_at(104.6, ly + 1.35, 'protected', 10, C_PROT, ha='left')

# ================= 좌측: One block =================
text_at(2, 50.5, 'One block', 14, INK, bold=True)
by = 44
for i, name in enumerate(['Attention', 'Merge', 'MLP']):
    bx = 2 + i * 11.5
    rrect(bx, by, 10, 4.6, fill=(BOXFM if name == 'Merge' else BOXF),
          line=(C_REG if name == 'Merge' else BOXE), line_w=1.3, round_amt=0.22)
    text_at(bx + 5, by + 2.3, name, 10.5, INK, bold=True, ha='center', va='center')
    if i < 2: line(bx + 10, by + 2.3, bx + 13, by + 2.3, color=ARROW, width=1.6, arrow=True)
text_at(36.5, by + 2.3, '×L', 13, INK, bold=True, italic=True, ha='left', va='center')
# Merge -> 메커니즘 콜아웃 점선(세로)
line(13, by, 13, 37.6, color=CALL, width=1.2, dash='sysDot')

# protect 인셋
text_at(2, 39, 'protect: pass through', 10, C_PROT, bold=True, ha='left')
row(3, 33.5, ['cls', 'reg', 'reg'], prot_n=3)
line(12.6, 34.8, 16.6, 34.8, color=ARROW, width=1.6, arrow=True)
row(18.0, 33.5, ['cls', 'reg', 'reg'], prot_n=3)

# merge 인셋
text_at(2, 30, 'merge r similar patches', 10, INK, bold=True, ha='left')
tok(4, 24.3, fc=C_PATCH); tok(8.4, 24.3, fc=C_PATCH)
text_at(7.25, 25.65, '≈', 12, MUT, ha='center', va='center')
line(11.5, 25.65, 16.5, 25.65, color=ARROW, width=1.6, arrow=True)
tok(17.3, 23.8, w=3.7, h=3.7, fc=C_PATCH, ec=MUT)
text_at(19.15, 22.0, 'size ↑', 8.5, MUT, ha='center')

# 좌/우 구분선
line(44, 3, 44, 52, color=DIVID, width=1.1, dash='dash')

# ================= 우측: Across blocks =================
text_at(47, 50.5, 'Across blocks', 14, INK, bold=True)
X1, X2, X3 = 54, 84, 110
def lane(y, name, color, snaps):
    text_at(46.5, y + 1.3, name, 12, color, bold=True, ha='left')
    ends = []
    for xs, (spec, pn) in zip([X1, X2, X3], snaps):
        ends.append(row(xs, y, spec, prot_n=pn))
    line(ends[0] + 1.5, y + 1.35, X2 - 1.5, y + 1.35, color=ARROW, width=1.6, arrow=True)
    line(ends[1] + 1.5, y + 1.35, X3 - 1.5, y + 1.35, color=ARROW, width=1.6, arrow=True)

tome = [(['cls'] + ['reg'] * 4 + ['patch'] * 3, 1),
        (['cls'] + ['reg'] * 2 + ['patch'] * 2, 1),
        (['cls'] + ['patch'] * 3,               1)]
ours = [(['cls'] + ['reg'] * 4 + ['patch'] * 3, 5),
        (['cls'] + ['reg'] * 4 + ['patch'] * 2, 5),
        (['cls'] + ['reg'] * 4 + ['patch'] * 1, 5)]
lane(37, 'ToMe', TOME_C, tome)
lane(16, 'Ours', REGBLU, ours)

# 첫 스냅샷 위 토큰 종류 브레이스 라벨
def brace_label(x0, x1, y, txt, color):
    line(x0, y, x0, y + 0.9, color=color, width=1.0)
    line(x0, y + 0.9, x1, y + 0.9, color=color, width=1.0)
    line(x1, y + 0.9, x1, y, color=color, width=1.0)
    text_at((x0 + x1) / 2, y + 2.0, txt, 9, color, ha='center')
brace_label(X1, X1 + 2.3, 40.5, 'CLS', C_CLS)
brace_label(X1 + 2.85, X1 + 2.85 + 4 * 2.85 - 0.55, 40.5, 'registers', C_REG)
brace_label(X1 + 2.85 + 4 * 2.85, X1 + 2.85 + 7 * 2.85 - 0.55, 40.5, 'patches', MUT)

# 깊이 라벨
for xx, lab in [(X1 + 3, 'block 1'), (X2 + 2, 'block L/2'), (X3 + 2, 'block L')]:
    text_at(xx, 11.5, lab, 9.5, GRAY66, ha='center')

prs.save(OUT)
import os
print("SAVED", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides._sldIdLst), "slide(s),",
      len(shp), "top-level shapes")
