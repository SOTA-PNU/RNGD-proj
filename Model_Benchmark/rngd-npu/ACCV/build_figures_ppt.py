#!/usr/bin/env python3
"""논문 4개 그림(teaser·주결과·PiToMe비교·ablation)을 PPT 네이티브 객체/차트로.
PNG 통삽입 아님 — 차트는 python-pptx CategoryChartData(데이터 pptx 내장), teaser는 도형(사각형/화살표).
색: Okabe-Ito(논문 그림과 통일) Ours=#0072B2, ToMe=#D55E00, PiToMe=#009E73. 출력 ACCV_그림_네이티브.pptx"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUT = "/home/jun/RNGD-proj/Model_Benchmark/rngd-npu/ACCV/ACCV_그림_네이티브.pptx"
OURS, TOME, PITOME = RGBColor(0x00, 0x72, 0xB2), RGBColor(0xD5, 0x55, 0x00), RGBColor(0x00, 0x9E, 0x73)
REG, PATCH, CLS, INK = OURS, RGBColor(0xD9, 0xDC, 0xE1), RGBColor(0x33, 0x33, 0x33), RGBColor(0x4D, 0x4D, 0x4D)
MUT = [RGBColor(0xCC, 0x79, 0xA7), RGBColor(0xE6, 0x9F, 0x00), RGBColor(0x56, 0xB4, 0xE9)]  # random/energy/highnorm

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0x14, 0x2a, 0x4a)
    return s


def txt(s, l, t, w, h, text, size=13, color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb


def box(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); sp.line.width = Pt(1)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def tokrow(s, x, y, keep_reg, n_patch=8, sq=0.3, gap=0.07):
    cx = x
    box(s, cx, y, sq, sq, CLS); cx += sq + gap                       # CLS
    for _ in range(4):
        box(s, cx, y, sq, sq, REG if keep_reg else PATCH); cx += sq + gap  # registers
    for _ in range(n_patch):
        box(s, cx, y, sq, sq, PATCH, line=RGBColor(0xC0, 0xC4, 0xCC)); cx += sq + gap
    return cx


def arrow(s, l, t, w, h, color, shape=MSO_SHAPE.RIGHT_ARROW, rot=0):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    sp.rotation = rot; sp.shadow.inherit = False
    return sp


def line_chart(s, cats, series, l, t, w, h, colors, ytitle, xtitle, smooth=False, markers=True):
    cd = CategoryChartData(); cd.categories = [str(c) for c in cats]
    for name, vals in series:
        cd.add_series(name, vals)
    ct = XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE
    gf = s.shapes.add_chart(ct, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = False
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(11)
    for i, ser in enumerate(ch.series):
        c = colors[i]
        ser.format.line.color.rgb = c; ser.format.line.width = Pt(2.5 if i == 0 else 2.0)
        ser.smooth = smooth
    va = ch.value_axis; va.has_title = True; va.axis_title.text_frame.text = ytitle
    va.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    va.tick_labels.font.size = Pt(10)
    ca = ch.category_axis; ca.has_title = True; ca.axis_title.text_frame.text = xtitle
    ca.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    ca.tick_labels.font.size = Pt(10)
    return gf


# =========================================================================
# 슬라이드 1 — TEASER (개념 도형)
# =========================================================================
s = slide("그림 1 (teaser) — 극단 압축에서 register 보호가 붕괴를 막는다")
txt(s, 0.6, 1.15, 3.5, 0.3, "DINOv2-reg 토큰", size=13, bold=True, color=INK)
tokrow(s, 0.6, 1.55, keep_reg=True)
# 범례 (경로 아래 빈 영역 — 분기 화살표와 겹치지 않게)
box(s, 0.9, 4.55, 0.26, 0.26, CLS); txt(s, 1.25, 4.52, 1.2, 0.3, "CLS", size=11)
box(s, 2.4, 4.55, 0.26, 0.26, REG); txt(s, 2.75, 4.52, 3.2, 0.3, "registers (전역 메모)", size=11)
box(s, 5.5, 4.55, 0.26, 0.26, PATCH, line=RGBColor(0xC0, 0xC4, 0xCC)); txt(s, 5.85, 4.52, 2.0, 0.3, "patches", size=11)
# 분기 화살표
arrow(s, 4.4, 1.4, 1.5, 0.5, TOME, rot=-18)   # 위로
arrow(s, 4.4, 2.1, 1.5, 0.5, OURS, rot=18)    # 아래로
# 위 경로: ToMe
txt(s, 6.2, 1.05, 6.0, 0.4, "ToMe: 유사도로 병합", size=14, bold=True, color=TOME)
tokrow(s, 6.2, 1.5, keep_reg=False, n_patch=2)   # register 사라짐
txt(s, 9.4, 1.35, 3.6, 0.35, "registers 합쳐져 소멸", size=11, italic=True, color=TOME)
txt(s, 9.4, 1.72, 3.6, 0.4, "kNN 63.99%  (붕괴)", size=15, bold=True, color=TOME)
# 아래 경로: Ours
txt(s, 6.2, 3.15, 6.5, 0.4, "Ours: register 보호 + patch만 병합", size=14, bold=True, color=OURS)
tokrow(s, 6.2, 3.6, keep_reg=True, n_patch=2)    # register 유지(파랑)
txt(s, 9.4, 3.45, 3.6, 0.35, "registers 유지", size=11, italic=True, color=OURS)
txt(s, 9.4, 3.82, 3.6, 0.4, "kNN 71.86%  (견고)", size=15, bold=True, color=OURS)
# +7.9 하이라이트
txt(s, 11.0, 2.5, 2.2, 0.6, "+7.9%p", size=26, bold=True, color=OURS)
# 캡션
txt(s, 0.6, 5.5, 12.1, 0.9,
    "92% 토큰 축소 · 학습 불필요 (DINOv2-reg ViT-B, ImageNet val).  "
    "register는 전역 정보를 담는데, 표준 병합은 극단 압축에서 이를 없애 버린다.",
    size=13, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# 슬라이드 2 — 주 결과 (Ours vs ToMe, 라인차트) + 음성대조군
# =========================================================================
s = slide("그림 2 — 주 결과: register 보호가 극단 압축서 정확도 보존")
comp = [37, 55, 74, 83, 92]
line_chart(s, comp,
           [("Ours (register 보호)", [75.68, 75.12, 74.20, 73.29, 71.86]),
            ("ToMe (CLS만)", [74.68, 72.95, 70.14, 67.79, 63.99])],
           0.7, 1.3, 7.2, 5.4, [OURS, TOME],
           "ImageNet kNN top-1 (%)", "토큰 축소 (%)")
txt(s, 8.2, 1.5, 4.6, 5.0,
    "• 무압축 76.33 (val LOO kNN).\n\n"
    "• 92% 축소: Ours 71.86 vs ToMe 63.99 → +7.87.\n\n"
    "• 격차는 압축 강해질수록 단조 증가.\n\n"
    "• 음성 대조군(register 없는 DINOv2): "
    "Ours=ToMe 라 Δ≈0(±0.05) — 이득 원천이 register임을 확인.",
    size=15, color=INK)

# =========================================================================
# 슬라이드 3 — 공식 PiToMe 비교 (정확도 vs FLOP + 처리량)
# =========================================================================
s = slide("그림 3 — 공식 PiToMe와 같은 예산 비교 (극단서 역전)")
flop = [17, 26, 35, 39, 43]
line_chart(s, flop,
           [("Ours", [75.68, 75.12, 74.20, 73.29, 71.86]),
            ("PiToMe", [76.07, 75.47, 73.84, 71.88, 68.70]),
            ("ToMe", [74.68, 72.95, 70.14, 67.79, 63.99])],
           0.7, 1.25, 6.1, 4.5, [OURS, PITOME, TOME],
           "kNN top-1 (%)", "FLOP 절감 (%)")
flopt = [0, 17, 26, 35, 39, 43]
line_chart(s, flopt,
           [("Ours", [349, 406, 452, 505, 545, 576]),
            ("PiToMe", [349, 395, 440, 493, 532, 563]),
            ("ToMe", [350, 405, 451, 504, 543, 574])],
           7.0, 1.25, 6.0, 4.5, [OURS, PITOME, TOME],
           "처리량 (im/s)", "FLOP 절감 (%)")
txt(s, 0.7, 6.0, 12.4, 1.0,
    "온건 압축(≤26% FLOP)선 PiToMe 우세 → 극단(≥35% FLOP)서 Ours 역전(교차점 토큰~74%, 92%서 +3.2).  "
    "처리량은 셋이 유사(Ours≈ToMe, PiToMe 약간 낮음) → 정확도 이득이 속도 대가 없이 온다.",
    size=13, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# 슬라이드 4 — ablation (5전략, 라인차트)
# =========================================================================
s = slide("그림 4 — keep-prior ablation: register만 효과")
compa = [37, 55, 74, 83, 91]
line_chart(s, compa,
           [("Ours (register)", [75.71, 75.30, 74.14, 73.25, 71.98]),
            ("ToMe (CLS만)", [74.70, 72.95, 70.36, 67.60, 63.90]),
            ("random", [74.70, 72.82, 70.00, 67.26, 63.19]),
            ("energy", [74.58, 72.70, 70.16, 67.51, 64.04]),
            ("high-norm", [74.78, 72.87, 70.16, 67.53, 63.84])],
           0.7, 1.3, 7.2, 5.4, [OURS, TOME] + MUT,
           "ImageNet kNN top-1 (%)", "토큰 축소 (%)")
txt(s, 8.2, 1.6, 4.6, 4.8,
    "같은 크기-가중 병합, 보호 대상만 교체(개수 동일).\n\n"
    "• register 보호만 baseline을 넘음.\n\n"
    "• random·energy·high-norm은 모두 무보호(ToMe) 노이즈 바닥에 뭉침.\n\n"
    "→ 이득은 '토큰을 더 보호해서'가 아니라 'register라서'.",
    size=15, color=INK)

prs.save(OUT)
print("SAVED", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides._sldIdLst), "slides")
